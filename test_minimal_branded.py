#!/usr/bin/env python3
"""
Minimal test to isolate PowerPoint corruption source
"""

import sys
sys.path.append('lib')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def test_minimal_branded():
    """Create minimal branded presentation without any complex components"""
    
    print("Creating minimal branded presentation...")
    
    # Start from blank presentation (no template)
    prs = Presentation()
    
    # Add title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Minimal branded test"
    
    # Add content slide with basic text
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    title.text = "Content Slide"
    content.text = "This is a simple content slide with no complex components"
    
    # Add blank slide with safe text box
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    text_box = slide3.shapes.add_textbox(
        Inches(1.0),  # Ensure positive coordinates
        Inches(1.0),
        Inches(8.0),  # Safe width
        Inches(1.0)   # Safe height
    )
    text_frame = text_box.text_frame
    text_frame.text = "Safe text box with validated coordinates"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(18)  # Safe font size
    p.font.bold = True
    
    # Save
    prs.save('test_minimal_branded.pptx')
    
    import os
    size = os.path.getsize('test_minimal_branded.pptx')
    print(f"Generated test_minimal_branded.pptx ({size:,} bytes)")
    print("This file should open without any corruption issues")

if __name__ == "__main__":
    test_minimal_branded()