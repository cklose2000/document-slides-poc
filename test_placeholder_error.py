#!/usr/bin/env python3
"""Test to isolate placeholder error"""

from pptx import Presentation
from pptx.util import Inches
import os

def test_placeholder_access():
    """Test accessing placeholders in different ways"""
    
    # Test 1: Create blank presentation
    print("Test 1: Blank presentation")
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Try to access placeholder_format
        for shape in slide.shapes:
            print(f"  Shape type: {type(shape)}")
            if hasattr(shape, 'placeholder_format'):
                print(f"  Has placeholder_format: {shape.placeholder_format}")
                if shape.placeholder_format:
                    print(f"  Placeholder type: {shape.placeholder_format.type}")
            else:
                print(f"  No placeholder_format attribute")
                
    except Exception as e:
        print(f"  Error: {e}")
        import traceback
        traceback.print_exc()
    
    # Test 2: Template presentation
    print("\nTest 2: Template presentation")
    template_path = "templates/firm_template.pptx"
    if os.path.exists(template_path):
        try:
            prs = Presentation(template_path)
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            
            for shape in slide.shapes:
                print(f"  Shape type: {type(shape)}")
                if hasattr(shape, 'placeholder_format'):
                    print(f"  Has placeholder_format: {shape.placeholder_format}")
                    if shape.placeholder_format:
                        print(f"  Placeholder type: {shape.placeholder_format.type}")
                else:
                    print(f"  No placeholder_format attribute")
                    
        except Exception as e:
            print(f"  Error: {e}")
            import traceback
            traceback.print_exc()
    else:
        print("  Template not found")
    
    # Test 3: Check if we can trigger the error
    print("\nTest 3: Trying to trigger 'shape is not a placeholder' error")
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Add a regular shape (not a placeholder)
        shape = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(2))
        
        # Try to access it as a placeholder (this might trigger the error)
        print(f"  Regular shape type: {type(shape)}")
        
        # This might be where the error comes from
        if hasattr(shape, 'is_placeholder'):
            print(f"  is_placeholder: {shape.is_placeholder}")
        
        # Try to clear a non-placeholder text frame
        if hasattr(shape, 'text_frame'):
            print(f"  Has text_frame: {shape.text_frame}")
            # Uncommenting this might trigger the error:
            # shape.text_frame.clear()
                    
    except Exception as e:
        print(f"  Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    print("🔍 Testing placeholder access patterns")
    print("=" * 50)
    test_placeholder_access()