"""
Create a sample branded PowerPoint template for testing
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_sample_template(output_path: str = 'templates/sample_brand_template.pptx'):
    """Create a sample branded template with custom colors and fonts"""
    
    # Create new presentation
    prs = Presentation()
    
    # Get slide master for customization
    slide_master = prs.slide_master
    
    # Define brand colors
    brand_primary = RGBColor(30, 70, 135)      # Deep blue
    brand_secondary = RGBColor(220, 130, 50)   # Orange accent
    brand_accent1 = RGBColor(100, 150, 200)    # Light blue
    brand_accent2 = RGBColor(50, 120, 80)      # Green
    
    # Create title slide
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    
    # Customize title slide
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Sample Brand Template"
    subtitle.text = "Professional Presentation Template"
    
    # Style title
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.name = 'Arial'
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.color.rgb = brand_primary
    title_paragraph.font.bold = True
    
    # Style subtitle  
    subtitle_paragraph = subtitle.text_frame.paragraphs[0]
    subtitle_paragraph.font.name = 'Arial'
    subtitle_paragraph.font.size = Pt(24)
    subtitle_paragraph.font.color.rgb = brand_secondary
    
    # Add brand rectangle
    brand_rect = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8), Inches(10), Inches(0.7)
    )
    brand_rect.fill.solid()
    brand_rect.fill.fore_color.rgb = brand_primary
    brand_rect.line.fill.background()
    
    # Create content slide layout
    content_layout = prs.slide_layouts[1]
    content_slide = prs.slides.add_slide(content_layout)
    
    # Customize content slide
    content_title = content_slide.shapes.title
    content_body = content_slide.placeholders[1]
    
    content_title.text = "Content Slide Example"
    content_body.text = "• Branded bullet point\n• Professional styling\n• Consistent colors"
    
    # Style content title
    content_title_paragraph = content_title.text_frame.paragraphs[0]
    content_title_paragraph.font.name = 'Arial'
    content_title_paragraph.font.size = Pt(32)
    content_title_paragraph.font.color.rgb = brand_primary
    content_title_paragraph.font.bold = True
    
    # Style content body
    for paragraph in content_body.text_frame.paragraphs:
        paragraph.font.name = 'Arial'
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(50, 50, 50)  # Dark gray
    
    # Add brand accent
    accent_rect = content_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.3), Inches(7.5)
    )
    accent_rect.fill.solid()
    accent_rect.fill.fore_color.rgb = brand_secondary
    accent_rect.line.fill.background()
    
    # Create blank slide for custom layouts
    blank_layout = prs.slide_layouts[6]
    blank_slide = prs.slides.add_slide(blank_layout)
    
    # Add sample chart slide
    chart_title = blank_slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(8), Inches(1)
    )
    chart_title.text_frame.text = "Financial Summary"
    
    # Style chart title
    chart_title_paragraph = chart_title.text_frame.paragraphs[0]
    chart_title_paragraph.font.name = 'Arial'
    chart_title_paragraph.font.size = Pt(32)
    chart_title_paragraph.font.color.rgb = brand_primary
    chart_title_paragraph.font.bold = True
    chart_title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add sample table with brand colors
    table_shape = blank_slide.shapes.add_table(4, 3, Inches(1), Inches(2), Inches(8), Inches(3))
    table = table_shape.table
    
    # Header row
    headers = ['Metric', 'Value', 'Source']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = brand_primary
        
        # White text for contrast
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True
                run.font.name = 'Arial'
                run.font.size = Pt(14)
    
    # Sample data rows
    sample_data = [
        ['Revenue', '$2.5M', 'Sheet1:B5'],
        ['Growth', '15%', 'Sheet1:C10'],
        ['Margin', '22%', 'Sheet1:D15']
    ]
    
    for row_idx, row_data in enumerate(sample_data, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
            
            # Style data cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Arial'
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(50, 50, 50)
    
    # Add footer with brand accent
    footer_rect = blank_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8), Inches(10), Inches(0.7)
    )
    footer_rect.fill.solid()
    footer_rect.fill.fore_color.rgb = brand_accent1
    footer_rect.line.fill.background()
    
    # Remove the sample slides (keep only layouts)
    # Note: We'll keep them for demo purposes
    
    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    # Save template
    prs.save(output_path)
    print(f"Sample brand template created: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_sample_template()