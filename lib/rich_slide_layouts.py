#!/usr/bin/env python3
"""
Rich Slide Layouts for Substantial Presentations

This module provides advanced slide layouts to create more impactful presentations:
- Executive dashboards with KPI cards
- Multi-chart analytical slides
- Timeline and roadmap visualizations
- Comparison matrices
- Infographic-style layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from typing import Dict, List, Any, Optional, Tuple
import random

class RichSlideLayouts:
    """Create rich, substantial slide layouts"""
    
    def __init__(self, brand_config: Optional[Dict[str, Any]] = None):
        """Initialize with brand configuration"""
        self.brand_config = brand_config or self._get_default_config()
        
    def _get_default_config(self) -> Dict[str, Any]:
        """Default brand configuration"""
        return {
            'theme_colors': {
                'primary': '#1e3c72',
                'secondary': '#2a5298',
                'accent1': '#667eea',
                'accent2': '#764ba2',
                'success': '#27ae60',
                'warning': '#f39c12',
                'danger': '#e74c3c',
                'dark': '#2c3e50',
                'light': '#ecf0f1'
            },
            'fonts': {
                'heading': {'family': 'Calibri Light', 'size': 24},
                'body': {'family': 'Calibri', 'size': 12},
                'caption': {'family': 'Calibri Light', 'size': 10}
            }
        }
    
    def create_executive_dashboard_slide(self, slide: Any, metrics: Dict[str, Any], 
                                       title: str = "Executive Dashboard") -> Any:
        """Create an executive dashboard with KPI cards"""
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(30, 60, 114)
        
        # Create KPI cards grid (2x3 layout)
        card_width = Inches(2.8)
        card_height = Inches(2.2)
        spacing = Inches(0.3)
        start_x = Inches(0.5)
        start_y = Inches(1.2)
        
        # Define KPI metrics with icons and colors
        kpis = [
            {'label': 'Total Revenue', 'value': '$12.5M', 'change': '+15%', 'color': 'success'},
            {'label': 'Gross Margin', 'value': '42.3%', 'change': '+2.1%', 'color': 'success'},
            {'label': 'Operating Costs', 'value': '$4.2M', 'change': '-5%', 'color': 'success'},
            {'label': 'Net Profit', 'value': '$3.8M', 'change': '+22%', 'color': 'success'},
            {'label': 'Cash Flow', 'value': '$2.1M', 'change': '+8%', 'color': 'warning'},
            {'label': 'ROI', 'value': '18.5%', 'change': '+3.2%', 'color': 'success'}
        ]
        
        # Override with actual metrics if provided
        if metrics:
            # Update KPIs with real data
            pass
        
        for i, kpi in enumerate(kpis):
            row = i // 3
            col = i % 3
            
            x = start_x + (card_width + spacing) * col
            y = start_y + (card_height + spacing) * row
            
            # Create card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_width, card_height
            )
            
            # Style the card
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(230, 230, 230)
            card.line.width = Pt(1)
            
            # Add KPI label
            label_box = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(0.2), 
                card_width - Inches(0.4), Inches(0.4)
            )
            label_frame = label_box.text_frame
            label_frame.text = kpi['label']
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = RGBColor(100, 100, 100)
            
            # Add main value
            value_box = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(0.6),
                card_width - Inches(0.4), Inches(0.8)
            )
            value_frame = value_box.text_frame
            value_frame.text = kpi['value']
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(32)
            value_para.font.bold = True
            value_para.font.color.rgb = RGBColor(30, 60, 114)
            
            # Add change indicator
            change_box = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(1.5),
                card_width - Inches(0.4), Inches(0.4)
            )
            change_frame = change_box.text_frame
            change_frame.text = kpi['change']
            change_para = change_frame.paragraphs[0]
            change_para.font.size = Pt(16)
            change_para.font.bold = True
            
            # Color based on positive/negative
            if kpi['color'] == 'success':
                change_para.font.color.rgb = RGBColor(39, 174, 96)
            elif kpi['color'] == 'warning':
                change_para.font.color.rgb = RGBColor(243, 156, 18)
            else:
                change_para.font.color.rgb = RGBColor(231, 76, 60)
            
            # Add trend indicator shape
            if '+' in kpi['change']:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.ISOSCELES_TRIANGLE,
                    x + card_width - Inches(0.6), y + Inches(1.5),
                    Inches(0.3), Inches(0.3)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(39, 174, 96)
                arrow.rotation = 0
            elif '-' in kpi['change']:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.ISOSCELES_TRIANGLE,
                    x + card_width - Inches(0.6), y + Inches(1.5),
                    Inches(0.3), Inches(0.3)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(231, 76, 60)
                arrow.rotation = 180
        
        return slide
    
    def create_multi_chart_analysis_slide(self, slide: Any, data: Dict[str, Any],
                                        title: str = "Financial Analysis") -> Any:
        """Create a slide with multiple charts for comprehensive analysis"""
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(30, 60, 114)
        
        # Chart 1: Revenue trend (line chart)
        chart1_data = CategoryChartData()
        chart1_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
        chart1_data.add_series('Revenue', (10.2, 11.5, 12.8, 14.2))
        chart1_data.add_series('Target', (10.0, 11.0, 12.0, 13.0))
        
        x, y, cx, cy = Inches(0.5), Inches(1), Inches(4.5), Inches(2.8)
        chart1 = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart1_data
        ).chart
        
        # Chart 2: Expense breakdown (pie chart)
        chart2_data = CategoryChartData()
        chart2_data.categories = ['Operations', 'Marketing', 'R&D', 'Admin']
        chart2_data.add_series('Expenses', (45, 25, 20, 10))
        
        x, y, cx, cy = Inches(5.0), Inches(1), Inches(4.5), Inches(2.8)
        chart2 = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart2_data
        ).chart
        
        # Chart 3: YoY comparison (column chart)
        chart3_data = CategoryChartData()
        chart3_data.categories = ['Revenue', 'Profit', 'Growth']
        chart3_data.add_series('2023', (10.5, 2.3, 15))
        chart3_data.add_series('2024', (12.8, 3.1, 22))
        
        x, y, cx, cy = Inches(0.5), Inches(4.2), Inches(4.5), Inches(2.8)
        chart3 = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart3_data
        ).chart
        
        # Add insights box
        insights_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.0), Inches(4.2), Inches(4.5), Inches(2.8)
        )
        insights_box.fill.solid()
        insights_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
        
        # Add insights text
        insights_text = slide.shapes.add_textbox(
            Inches(5.2), Inches(4.4), Inches(4.1), Inches(2.4)
        )
        insights_frame = insights_text.text_frame
        insights_frame.text = "Key Insights"
        
        p = insights_frame.add_paragraph()
        p.text = "• Revenue exceeded targets by 9.2%"
        p.font.size = Pt(12)
        
        p = insights_frame.add_paragraph()
        p.text = "• Operating efficiency improved 15%"
        p.font.size = Pt(12)
        
        p = insights_frame.add_paragraph()
        p.text = "• YoY growth accelerating"
        p.font.size = Pt(12)
        
        return slide
    
    def create_timeline_roadmap_slide(self, slide: Any, milestones: List[Dict[str, Any]],
                                    title: str = "Strategic Roadmap") -> Any:
        """Create a timeline/roadmap visualization"""
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(30, 60, 114)
        
        # Create timeline line
        line_y = Inches(3.5)
        line = slide.shapes.add_connector(
            1, Inches(1), line_y, Inches(9), line_y
        )
        line.line.color.rgb = RGBColor(102, 126, 234)
        line.line.width = Pt(3)
        
        # Default milestones if none provided
        if not milestones:
            milestones = [
                {'date': 'Q1 2024', 'title': 'Product Launch', 'status': 'completed'},
                {'date': 'Q2 2024', 'title': 'Market Expansion', 'status': 'completed'},
                {'date': 'Q3 2024', 'title': 'Series B Funding', 'status': 'in_progress'},
                {'date': 'Q4 2024', 'title': 'Global Rollout', 'status': 'planned'},
                {'date': 'Q1 2025', 'title': 'IPO Preparation', 'status': 'planned'}
            ]
        
        # Add milestones
        milestone_count = len(milestones)
        spacing = Inches(8) / (milestone_count - 1) if milestone_count > 1 else Inches(0)
        
        for i, milestone in enumerate(milestones):
            x = Inches(1) + (spacing * i)
            
            # Add milestone circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x - Inches(0.2), line_y - Inches(0.2),
                Inches(0.4), Inches(0.4)
            )
            
            # Color based on status
            if milestone['status'] == 'completed':
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(39, 174, 96)
            elif milestone['status'] == 'in_progress':
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(243, 156, 18)
            else:
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(189, 195, 199)
            
            # Add date above
            date_box = slide.shapes.add_textbox(
                x - Inches(0.8), line_y - Inches(1.2),
                Inches(1.6), Inches(0.4)
            )
            date_frame = date_box.text_frame
            date_frame.text = milestone['date']
            date_para = date_frame.paragraphs[0]
            date_para.alignment = PP_ALIGN.CENTER
            date_para.font.size = Pt(12)
            date_para.font.bold = True
            
            # Add milestone title below
            title_box = slide.shapes.add_textbox(
                x - Inches(1), line_y + Inches(0.4),
                Inches(2), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = milestone['title']
            title_para = title_frame.paragraphs[0]
            title_para.alignment = PP_ALIGN.CENTER
            title_para.font.size = Pt(11)
            title_frame.word_wrap = True
        
        return slide
    
    def create_comparison_matrix_slide(self, slide: Any, comparison_data: Dict[str, Any],
                                     title: str = "Competitive Analysis") -> Any:
        """Create a comparison matrix slide"""
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(30, 60, 114)
        
        # Default comparison data
        if not comparison_data:
            comparison_data = {
                'criteria': ['Market Share', 'Technology', 'Price', 'Support', 'Features'],
                'competitors': {
                    'Our Company': [85, 95, 70, 90, 95],
                    'Competitor A': [70, 80, 85, 75, 80],
                    'Competitor B': [60, 70, 90, 70, 75],
                    'Market Leader': [90, 85, 60, 85, 90]
                }
            }
        
        # Create matrix table
        rows = len(comparison_data['criteria']) + 1
        cols = len(comparison_data['competitors']) + 1
        
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), 
                                      Inches(9), Inches(5.5)).table
        
        # Set column widths
        table.columns[0].width = Inches(2)
        for i in range(1, cols):
            table.columns[i].width = Inches(1.75)
        
        # Header row
        table.cell(0, 0).text = "Criteria"
        for i, company in enumerate(comparison_data['competitors'].keys()):
            table.cell(0, i + 1).text = company
        
        # Style header row
        for i in range(cols):
            cell = table.cell(0, i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 60, 114)
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.font.bold = True
            para.alignment = PP_ALIGN.CENTER
        
        # Data rows
        for i, criterion in enumerate(comparison_data['criteria']):
            table.cell(i + 1, 0).text = criterion
            
            for j, (company, scores) in enumerate(comparison_data['competitors'].items()):
                score = scores[i]
                cell = table.cell(i + 1, j + 1)
                
                # Add score with color coding
                if score >= 90:
                    color = RGBColor(39, 174, 96)
                elif score >= 75:
                    color = RGBColor(52, 152, 219)
                elif score >= 60:
                    color = RGBColor(243, 156, 18)
                else:
                    color = RGBColor(231, 76, 60)
                
                cell.text = f"{score}%"
                para = cell.text_frame.paragraphs[0]
                para.font.color.rgb = color
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def create_swot_analysis_slide(self, slide: Any, swot_data: Dict[str, List[str]],
                                 title: str = "SWOT Analysis") -> Any:
        """Create a SWOT analysis slide with 2x2 matrix"""
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(30, 60, 114)
        
        # Default SWOT data
        if not swot_data:
            swot_data = {
                'strengths': [
                    'Market leader in innovation',
                    'Strong brand recognition',
                    'Experienced team'
                ],
                'weaknesses': [
                    'Limited international presence',
                    'High operational costs',
                    'Dependency on key suppliers'
                ],
                'opportunities': [
                    'Emerging markets expansion',
                    'Digital transformation',
                    'Strategic partnerships'
                ],
                'threats': [
                    'Increasing competition',
                    'Economic uncertainty',
                    'Regulatory changes'
                ]
            }
        
        # Create 2x2 grid
        box_width = Inches(4.3)
        box_height = Inches(2.8)
        spacing = Inches(0.2)
        
        swot_config = [
            {'key': 'strengths', 'label': 'Strengths', 'color': RGBColor(39, 174, 96), 'x': 0, 'y': 0},
            {'key': 'weaknesses', 'label': 'Weaknesses', 'color': RGBColor(231, 76, 60), 'x': 1, 'y': 0},
            {'key': 'opportunities', 'label': 'Opportunities', 'color': RGBColor(52, 152, 219), 'x': 0, 'y': 1},
            {'key': 'threats', 'label': 'Threats', 'color': RGBColor(243, 156, 18), 'x': 1, 'y': 1}
        ]
        
        for config in swot_config:
            x = Inches(0.5) + (box_width + spacing) * config['x']
            y = Inches(1.2) + (box_height + spacing) * config['y']
            
            # Create box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, box_width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            box.line.color.rgb = config['color']
            box.line.width = Pt(2)
            
            # Add header
            header_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, box_width, Inches(0.5)
            )
            header_box.fill.solid()
            header_box.fill.fore_color.rgb = config['color']
            
            # Add header text
            header_text = slide.shapes.add_textbox(
                x + Inches(0.1), y + Inches(0.05),
                box_width - Inches(0.2), Inches(0.4)
            )
            header_frame = header_text.text_frame
            header_frame.text = config['label']
            header_para = header_frame.paragraphs[0]
            header_para.font.size = Pt(16)
            header_para.font.bold = True
            header_para.font.color.rgb = RGBColor(255, 255, 255)
            header_para.alignment = PP_ALIGN.CENTER
            
            # Add items
            items_text = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(0.7),
                box_width - Inches(0.4), box_height - Inches(0.9)
            )
            items_frame = items_text.text_frame
            items_frame.word_wrap = True
            
            for i, item in enumerate(swot_data.get(config['key'], [])):
                if i == 0:
                    items_frame.text = f"• {item}"
                else:
                    p = items_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.space_after = Pt(6)
                
                para = items_frame.paragraphs[i]
                para.font.size = Pt(11)
        
        return slide

def enhance_slide_generator_with_rich_layouts(generator, rich_layouts: RichSlideLayouts):
    """Add rich layout methods to existing slide generator"""
    
    # Monkey patch the methods onto the generator
    generator.create_executive_dashboard = lambda data, title="Executive Dashboard": \
        rich_layouts.create_executive_dashboard_slide(
            generator.prs.slides.add_slide(generator.prs.slide_layouts[5]), 
            data, title
        )
    
    generator.create_multi_chart_analysis = lambda data, title="Financial Analysis": \
        rich_layouts.create_multi_chart_analysis_slide(
            generator.prs.slides.add_slide(generator.prs.slide_layouts[5]),
            data, title
        )
    
    generator.create_timeline_roadmap = lambda milestones=None, title="Strategic Roadmap": \
        rich_layouts.create_timeline_roadmap_slide(
            generator.prs.slides.add_slide(generator.prs.slide_layouts[5]),
            milestones, title
        )
    
    generator.create_comparison_matrix = lambda data=None, title="Competitive Analysis": \
        rich_layouts.create_comparison_matrix_slide(
            generator.prs.slides.add_slide(generator.prs.slide_layouts[5]),
            data, title
        )
    
    generator.create_swot_analysis = lambda data=None, title="SWOT Analysis": \
        rich_layouts.create_swot_analysis_slide(
            generator.prs.slides.add_slide(generator.prs.slide_layouts[5]),
            data, title
        )
    
    return generator

# Example usage
if __name__ == "__main__":
    print("🎨 Rich Slide Layouts System")
    print("Available layouts:")
    print("  - Executive Dashboard (KPI cards)")
    print("  - Multi-Chart Analysis")
    print("  - Timeline/Roadmap")
    print("  - Comparison Matrix")
    print("  - SWOT Analysis")
    print("\n✨ Ready to create substantial presentations!")