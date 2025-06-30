"""
Reusable slide components for professional presentations.
Provides modular components that can be composed to create various slide types.
"""

from typing import Dict, List, Tuple, Optional, Any, Union
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import PicturePlaceholder
from datetime import datetime
import math


class SlideComponentBase:
    """Base class for all slide components."""
    
    def __init__(self, slide: Slide):
        self.slide = slide
        self.shapes = slide.shapes
        
    def _get_brand_colors(self) -> Dict[str, RGBColor]:
        """Get standard brand colors for consistency."""
        return {
            'primary': RGBColor(0, 47, 108),      # Navy blue
            'secondary': RGBColor(0, 176, 240),   # Light blue
            'accent': RGBColor(255, 192, 0),      # Gold
            'success': RGBColor(34, 139, 34),     # Green
            'warning': RGBColor(255, 140, 0),     # Orange
            'danger': RGBColor(220, 20, 60),      # Red
            'neutral': RGBColor(128, 128, 128),   # Gray
            'dark': RGBColor(51, 51, 51),         # Dark gray
            'light': RGBColor(242, 242, 242)      # Light gray
        }
    
    def _format_number(self, value: Union[int, float], decimal_places: int = 0) -> str:
        """Format number with appropriate suffix (K, M, B)."""
        if value == 0:
            return "0"
        
        abs_value = abs(value)
        sign = "-" if value < 0 else ""
        
        if abs_value >= 1e9:
            return f"{sign}{abs_value/1e9:.{decimal_places}f}B"
        elif abs_value >= 1e6:
            return f"{sign}{abs_value/1e6:.{decimal_places}f}M"
        elif abs_value >= 1e3:
            return f"{sign}{abs_value/1e3:.{decimal_places}f}K"
        else:
            return f"{sign}{abs_value:.{decimal_places}f}"
    
    def _format_percentage(self, value: float, decimal_places: int = 1) -> str:
        """Format percentage with + or - sign."""
        sign = "+" if value > 0 else ""
        return f"{sign}{value:.{decimal_places}f}%"


class TextComponents(SlideComponentBase):
    """Components for text elements on slides."""
    
    def add_title(self, text: str, left: float, top: float, width: float, height: float,
                  font_size: int = 28, bold: bool = True) -> BaseShape:
        """Add a title text box."""
        text_box = self.shapes.add_textbox(Inches(left), Inches(top), 
                                           Inches(width), Inches(height))
        text_frame = text_box.text_frame
        text_frame.clear()
        
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = self._get_brand_colors()['dark']
        p.alignment = PP_ALIGN.LEFT
        
        text_frame.margin_left = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.word_wrap = True
        
        return text_box
    
    def add_body_text(self, text: str, left: float, top: float, width: float, height: float,
                      font_size: int = 12, line_spacing: float = 1.2) -> BaseShape:
        """Add body text with proper formatting."""
        text_box = self.shapes.add_textbox(Inches(left), Inches(top), 
                                           Inches(width), Inches(height))
        text_frame = text_box.text_frame
        text_frame.clear()
        
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = self._get_brand_colors()['dark']
        p.line_spacing = line_spacing
        
        text_frame.word_wrap = True
        
        return text_box
    
    def add_bullet_points(self, items: List[str], left: float, top: float, 
                         width: float, height: float, font_size: int = 11,
                         bullet_char: str = "•") -> BaseShape:
        """Add formatted bullet points."""
        text_box = self.shapes.add_textbox(Inches(left), Inches(top), 
                                           Inches(width), Inches(height))
        text_frame = text_box.text_frame
        text_frame.clear()
        
        for i, item in enumerate(items):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = f"{bullet_char} {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = self._get_brand_colors()['dark']
            p.level = 0
            p.line_spacing = 1.2
        
        text_frame.word_wrap = True
        
        return text_box


class DataComponents(SlideComponentBase):
    """Components for data visualization elements."""
    
    def add_kpi_card(self, title: str, value: Union[str, float], change: Optional[float],
                     left: float, top: float, width: float = 2.0, height: float = 1.5,
                     data_point_id: Optional[str] = None) -> BaseShape:
        """Add a KPI card with value and change indicator."""
        # Card background
        card = self.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                     Inches(left), Inches(top), 
                                     Inches(width), Inches(height))
        
        # Style the card
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = self._get_brand_colors()['light']
        card.line.width = Pt(1)
        
        # Add shadow effect
        card.shadow.visible = True
        card.shadow.distance = Pt(2)
        card.shadow.blur_radius = Pt(4)
        card.shadow.transparency = 0.8
        
        # Title
        title_box = self.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.1),
                                            Inches(width - 0.2), Inches(0.3))
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.add_paragraph()
        p.text = title.upper()
        p.font.size = Pt(10)
        p.font.color.rgb = self._get_brand_colors()['neutral']
        p.font.bold = True
        
        # Value
        value_text = self._format_number(value, 1) if isinstance(value, (int, float)) else str(value)
        value_box = self.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.4),
                                            Inches(width - 0.2), Inches(0.6))
        value_frame = value_box.text_frame
        value_frame.clear()
        p = value_frame.add_paragraph()
        p.text = value_text
        p.font.size = Pt(24)
        p.font.color.rgb = self._get_brand_colors()['primary']
        p.font.bold = True
        
        # Change indicator
        if change is not None:
            change_color = self._get_brand_colors()['success'] if change > 0 else self._get_brand_colors()['danger']
            arrow = "▲" if change > 0 else "▼"
            
            change_box = self.shapes.add_textbox(Inches(left + 0.1), Inches(top + 1.0),
                                                 Inches(width - 0.2), Inches(0.3))
            change_frame = change_box.text_frame
            change_frame.clear()
            p = change_frame.add_paragraph()
            p.text = f"{arrow} {self._format_percentage(change)}"
            p.font.size = Pt(12)
            p.font.color.rgb = change_color
            p.font.bold = True
        
        # Add data point ID if provided
        if data_point_id:
            card.click_action.hyperlink.address = f"#source:{data_point_id}"
        
        return card
    
    def add_metric_table(self, data: List[Dict[str, Any]], left: float, top: float,
                        width: float, height: float, show_variance: bool = True) -> BaseShape:
        """Add a formatted metric comparison table."""
        rows = len(data) + 1  # +1 for header
        cols = len(data[0].keys()) if data else 3
        
        table = self.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                      Inches(width), Inches(height)).table
        
        # Style the table
        for row in table.rows:
            row.height = Inches(height / rows)
        
        # Add headers
        if data:
            headers = list(data[0].keys())
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._get_brand_colors()['primary']
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Add data rows
        for row_idx, row_data in enumerate(data):
            for col_idx, (key, value) in enumerate(row_data.items()):
                cell = table.cell(row_idx + 1, col_idx)
                
                # Format value based on type
                if isinstance(value, (int, float)):
                    if 'variance' in key.lower() or 'change' in key.lower():
                        cell.text = self._format_percentage(value)
                        # Color code variance
                        color = self._get_brand_colors()['success'] if value > 0 else self._get_brand_colors()['danger']
                        cell.text_frame.paragraphs[0].font.color.rgb = color
                    else:
                        cell.text = self._format_number(value, 1)
                else:
                    cell.text = str(value)
                
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                
                # Alternate row coloring
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._get_brand_colors()['light']
        
        return table
    
    def add_progress_indicator(self, progress: float, left: float, top: float,
                              width: float = 3.0, height: float = 0.3,
                              label: Optional[str] = None) -> BaseShape:
        """Add a progress bar indicator."""
        # Background bar
        bg_bar = self.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = self._get_brand_colors()['light']
        bg_bar.line.fill.background()
        
        # Progress bar
        progress_width = width * (progress / 100)
        if progress_width > 0:
            progress_bar = self.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                Inches(left), Inches(top),
                                                Inches(progress_width), Inches(height))
            progress_bar.fill.solid()
            
            # Color based on progress
            if progress >= 75:
                color = self._get_brand_colors()['success']
            elif progress >= 50:
                color = self._get_brand_colors()['warning']
            else:
                color = self._get_brand_colors()['danger']
            
            progress_bar.fill.fore_color.rgb = color
            progress_bar.line.fill.background()
        
        # Add label if provided
        if label:
            label_box = self.shapes.add_textbox(Inches(left + width + 0.1), 
                                               Inches(top - 0.05),
                                               Inches(1.0), Inches(height + 0.1))
            label_frame = label_box.text_frame
            label_frame.clear()
            p = label_frame.add_paragraph()
            p.text = f"{label}: {progress:.0f}%"
            p.font.size = Pt(10)
            p.font.color.rgb = self._get_brand_colors()['dark']
            label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        return bg_bar


class VisualComponents(SlideComponentBase):
    """Components for visual elements like icons, shapes, and graphics."""
    
    def add_icon_with_text(self, icon_shape: int, text: str, left: float, top: float,
                          size: float = 0.5, color: Optional[RGBColor] = None) -> BaseShape:
        """Add an icon shape with accompanying text."""
        # Add icon
        icon = self.shapes.add_shape(icon_shape, Inches(left), Inches(top),
                                    Inches(size), Inches(size))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color or self._get_brand_colors()['primary']
        icon.line.fill.background()
        
        # Add text
        text_box = self.shapes.add_textbox(Inches(left + size + 0.1), Inches(top),
                                           Inches(2.0), Inches(size))
        text_frame = text_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = self._get_brand_colors()['dark']
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        return icon
    
    def add_divider(self, left: float, top: float, width: float,
                   thickness: float = 1, color: Optional[RGBColor] = None) -> BaseShape:
        """Add a horizontal divider line."""
        line = self.shapes.add_connector(1, Inches(left), Inches(top),
                                        Inches(left + width), Inches(top))
        line.line.width = Pt(thickness)
        line.line.color.rgb = color or self._get_brand_colors()['light']
        
        return line
    
    def add_callout_box(self, text: str, left: float, top: float,
                       width: float, height: float, style: str = 'info') -> BaseShape:
        """Add a callout box with styled background."""
        # Color mapping
        style_colors = {
            'info': self._get_brand_colors()['secondary'],
            'success': self._get_brand_colors()['success'],
            'warning': self._get_brand_colors()['warning'],
            'danger': self._get_brand_colors()['danger'],
            'neutral': self._get_brand_colors()['neutral']
        }
        
        # Background shape
        box = self.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = style_colors.get(style, style_colors['info'])
        box.fill.transparency = 0.9  # Make it subtle
        box.line.color.rgb = style_colors.get(style, style_colors['info'])
        box.line.width = Pt(2)
        
        # Text
        text_box = self.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.1),
                                           Inches(width - 0.2), Inches(height - 0.2))
        text_frame = text_box.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = self._get_brand_colors()['dark']
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        return box


class LayoutComponents(SlideComponentBase):
    """Components for layout structures and grids."""
    
    def create_grid_layout(self, rows: int, cols: int, 
                          left: float = 0.5, top: float = 1.5,
                          width: float = 9.0, height: float = 5.0,
                          padding: float = 0.1) -> List[Dict[str, float]]:
        """Create a grid layout and return cell positions."""
        cells = []
        cell_width = (width - (cols - 1) * padding) / cols
        cell_height = (height - (rows - 1) * padding) / rows
        
        for row in range(rows):
            for col in range(cols):
                cell_left = left + col * (cell_width + padding)
                cell_top = top + row * (cell_height + padding)
                cells.append({
                    'left': cell_left,
                    'top': cell_top,
                    'width': cell_width,
                    'height': cell_height
                })
        
        return cells
    
    def create_header_content_layout(self, header_height: float = 1.0) -> Dict[str, Dict[str, float]]:
        """Create a layout with header and content areas."""
        return {
            'header': {
                'left': 0.5,
                'top': 0.5,
                'width': 9.0,
                'height': header_height
            },
            'content': {
                'left': 0.5,
                'top': 0.5 + header_height + 0.2,
                'width': 9.0,
                'height': 6.5 - header_height - 0.7
            }
        }
    
    def create_sidebar_layout(self, sidebar_width: float = 3.0) -> Dict[str, Dict[str, float]]:
        """Create a layout with sidebar and main content."""
        return {
            'sidebar': {
                'left': 0.5,
                'top': 1.5,
                'width': sidebar_width,
                'height': 5.0
            },
            'main': {
                'left': 0.5 + sidebar_width + 0.3,
                'top': 1.5,
                'width': 9.0 - sidebar_width - 0.3,
                'height': 5.0
            }
        }
    
    def add_section_header(self, text: str, left: float, top: float, width: float) -> BaseShape:
        """Add a section header with underline."""
        # Header text
        header = TextComponents(self.slide).add_title(text, left, top, width, 0.4,
                                                     font_size=16, bold=True)
        
        # Underline
        line = VisualComponents(self.slide).add_divider(left, top + 0.4, width,
                                                       thickness=2,
                                                       color=self._get_brand_colors()['primary'])
        
        return header


class CompositeComponents(SlideComponentBase):
    """Higher-level components composed of multiple basic components."""
    
    def __init__(self, slide: Slide):
        super().__init__(slide)
        self.text = TextComponents(slide)
        self.data = DataComponents(slide)
        self.visual = VisualComponents(slide)
        self.layout = LayoutComponents(slide)
    
    def add_kpi_dashboard(self, kpis: List[Dict[str, Any]], left: float = 0.5, 
                         top: float = 1.5, arrangement: str = 'horizontal') -> List[BaseShape]:
        """Add a dashboard of KPI cards."""
        shapes = []
        
        if arrangement == 'horizontal':
            card_width = 2.0
            spacing = 0.3
            for i, kpi in enumerate(kpis):
                card_left = left + i * (card_width + spacing)
                shape = self.data.add_kpi_card(
                    title=kpi.get('title', 'KPI'),
                    value=kpi.get('value', 0),
                    change=kpi.get('change'),
                    left=card_left,
                    top=top,
                    width=card_width,
                    data_point_id=kpi.get('data_point_id')
                )
                shapes.append(shape)
        else:  # vertical
            card_height = 1.5
            spacing = 0.2
            for i, kpi in enumerate(kpis):
                card_top = top + i * (card_height + spacing)
                shape = self.data.add_kpi_card(
                    title=kpi.get('title', 'KPI'),
                    value=kpi.get('value', 0),
                    change=kpi.get('change'),
                    left=left,
                    top=card_top,
                    data_point_id=kpi.get('data_point_id')
                )
                shapes.append(shape)
        
        return shapes
    
    def add_comparison_matrix(self, data: Dict[str, Dict[str, Any]], 
                             left: float, top: float, width: float, height: float) -> BaseShape:
        """Add a comparison matrix with multiple entities and metrics."""
        # Convert dict to list format for table
        table_data = []
        entities = list(data.keys())
        
        if entities:
            metrics = list(data[entities[0]].keys())
            
            # Create header row
            headers = [''] + entities
            
            # Create data rows
            for metric in metrics:
                row = {'Metric': metric}
                for entity in entities:
                    row[entity] = data[entity].get(metric, 'N/A')
                table_data.append(row)
        
        return self.data.add_metric_table(table_data, left, top, width, height)
    
    def add_timeline_visualization(self, events: List[Dict[str, Any]], 
                                  left: float, top: float, width: float, height: float) -> List[BaseShape]:
        """Add a timeline visualization with milestones."""
        shapes = []
        
        # Timeline line
        timeline_y = top + height / 2
        timeline = self.visual.add_divider(left, timeline_y, width, thickness=3,
                                          color=self._get_brand_colors()['primary'])
        shapes.append(timeline)
        
        # Add events
        num_events = len(events)
        if num_events > 0:
            spacing = width / (num_events + 1)
            
            for i, event in enumerate(events):
                event_x = left + (i + 1) * spacing
                
                # Milestone circle
                circle = self.shapes.add_shape(MSO_SHAPE.OVAL,
                                              Inches(event_x - 0.15), Inches(timeline_y - 0.15),
                                              Inches(0.3), Inches(0.3))
                circle.fill.solid()
                circle.fill.fore_color.rgb = self._get_brand_colors()['accent']
                circle.line.color.rgb = self._get_brand_colors()['primary']
                circle.line.width = Pt(2)
                shapes.append(circle)
                
                # Event text (alternating above/below)
                text_y = timeline_y - 0.8 if i % 2 == 0 else timeline_y + 0.4
                text_box = self.text.add_body_text(
                    f"{event.get('date', '')}\n{event.get('title', '')}",
                    event_x - 0.5, text_y, 1.0, 0.5, font_size=9
                )
                shapes.append(text_box)
        
        return shapes
    
    def add_risk_matrix(self, risks: List[Dict[str, Any]], 
                       left: float, top: float, size: float = 4.0) -> List[BaseShape]:
        """Add a risk assessment matrix (impact vs probability)."""
        shapes = []
        
        # Matrix background
        matrix = self.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(left), Inches(top),
                                      Inches(size), Inches(size))
        matrix.fill.solid()
        matrix.fill.fore_color.rgb = self._get_brand_colors()['light']
        matrix.line.color.rgb = self._get_brand_colors()['neutral']
        shapes.append(matrix)
        
        # Grid lines
        for i in range(1, 3):
            # Horizontal
            h_line = self.visual.add_divider(left, top + i * size/3, size)
            shapes.append(h_line)
            # Vertical
            v_line = self.shapes.add_connector(1, 
                                              Inches(left + i * size/3), Inches(top),
                                              Inches(left + i * size/3), Inches(top + size))
            v_line.line.color.rgb = self._get_brand_colors()['neutral']
            shapes.append(v_line)
        
        # Labels
        impact_labels = ['Low', 'Medium', 'High']
        prob_labels = ['Low', 'Medium', 'High']
        
        # Y-axis label (Impact)
        impact_label = self.text.add_body_text('Impact →', left - 0.7, top + size/2 - 0.2, 
                                              0.6, 0.4, font_size=10)
        shapes.append(impact_label)
        
        # X-axis label (Probability)
        prob_label = self.text.add_body_text('Probability →', left + size/2 - 0.5, 
                                            top + size + 0.1, 1.0, 0.3, font_size=10)
        shapes.append(prob_label)
        
        # Plot risks
        for risk in risks:
            impact = risk.get('impact', 1)  # 1-3
            probability = risk.get('probability', 1)  # 1-3
            
            # Calculate position
            risk_x = left + (probability - 0.5) * size/3
            risk_y = top + size - (impact - 0.5) * size/3
            
            # Risk bubble
            bubble = self.shapes.add_shape(MSO_SHAPE.OVAL,
                                          Inches(risk_x - 0.2), Inches(risk_y - 0.2),
                                          Inches(0.4), Inches(0.4))
            
            # Color based on risk level
            risk_score = impact * probability
            if risk_score >= 6:
                color = self._get_brand_colors()['danger']
            elif risk_score >= 4:
                color = self._get_brand_colors()['warning']
            else:
                color = self._get_brand_colors()['success']
            
            bubble.fill.solid()
            bubble.fill.fore_color.rgb = color
            bubble.line.fill.background()
            shapes.append(bubble)
            
            # Risk label
            label = self.text.add_body_text(risk.get('name', 'Risk')[:10],
                                           risk_x - 0.3, risk_y + 0.25,
                                           0.6, 0.3, font_size=8)
            shapes.append(label)
        
        return shapes