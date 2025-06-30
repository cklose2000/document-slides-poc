"""
Brand-Aware Slide Generator

Enhanced slide generator that applies brand templates and styling
for consistent, professional presentation output.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
# from pptx.oxml.shared import qn  # Not needed for current implementation
try:
    from .template_parser import BrandManager, TemplateParser
    from .source_tracker import SourceTracker
    from .chart_generator import ChartGenerator
    from .visual_effects import VisualEffectsEngine
    from .simple_visual_effects import SimpleVisualEffects, enhance_slide_simply
    from .rich_slide_layouts import RichSlideLayouts
except ImportError:
    from template_parser import BrandManager, TemplateParser
    from source_tracker import SourceTracker
    from chart_generator import ChartGenerator
    from visual_effects import VisualEffectsEngine
    from simple_visual_effects import SimpleVisualEffects, enhance_slide_simply
    from rich_slide_layouts import RichSlideLayouts
import os
import re
from typing import Dict, List, Any, Optional, Tuple
from io import BytesIO

class BrandedSlideGenerator:
    """Generate slides with consistent brand styling from templates"""
    
    def __init__(self, brand_manager: BrandManager = None, template_name: str = None, 
                 source_tracker: Optional[SourceTracker] = None):
        """Initialize with brand manager, template selection, and source tracker"""
        self.brand_manager = brand_manager or BrandManager()
        self.source_tracker = source_tracker
        
        if template_name:
            self.brand_manager.set_current_template(template_name)
        
        # Get current brand configuration
        self.brand_config = self.brand_manager.get_current_brand_config()
        
        # Initialize chart generator with adapted brand config
        chart_config = self._adapt_brand_config_for_charts(self.brand_config)
        self.chart_generator = ChartGenerator(chart_config)
        
        # Initialize visual effects engines for modern styling
        self.visual_effects = VisualEffectsEngine(self.brand_config)
        self.simple_effects = SimpleVisualEffects(self.brand_config)
        
        # Initialize rich layouts for substantial slides
        self.rich_layouts = RichSlideLayouts(self.brand_config)
        
        # Initialize presentation
        self._init_presentation()
    
    def _init_presentation(self):
        """Initialize presentation with template or create blank"""
        current_template = self.brand_manager.get_current_template()
        
        if current_template and os.path.exists(current_template.template_path):
            # Start with template presentation
            self.prs = Presentation(current_template.template_path)
            # Note: We'll keep existing slides in template for now to avoid clearing issues
            # In production, we'd implement proper slide clearing
        else:
            # Create blank presentation
            self.prs = Presentation()
    
    def create_financial_summary_slide(self, data: Dict[str, Any], source_refs: Dict[str, Any]) -> Any:
        """Create a branded financial summary slide with modern visual effects"""
        # Select appropriate layout
        layout = self._get_layout_for_content('financial_summary')
        slide = self.prs.slides.add_slide(layout)
        
        # Apply simple visual effects for content slides
        enhance_slide_simply(slide, 'content', self.brand_config)
        
        # Add title with brand styling
        title_shape = self._add_branded_title(slide, "Financial Summary")
        
        # Add metrics table with brand colors
        if data and isinstance(data, dict):
            self._add_branded_metrics_table(slide, data, Inches(1), Inches(2))
        
        # Add source attribution
        self.add_source_attribution(slide, source_refs)
        
        return slide
    
    def create_company_overview_slide(self, company_data: Dict[str, Any], source_refs: Dict[str, Any]) -> Any:
        """Create a branded company overview slide with modern visual effects"""
        layout = self._get_layout_for_content('company_overview')
        slide = self.prs.slides.add_slide(layout)
        
        # Apply simple visual effects for content slides
        enhance_slide_simply(slide, 'content', self.brand_config)
        
        # Add title
        title_shape = self._add_branded_title(slide, "Company Overview")
        
        # Add company info with brand styling
        if company_data:
            self._add_branded_company_info(slide, company_data, Inches(1), Inches(2))
        
        # Add source attribution
        self.add_source_attribution(slide, source_refs)
        
        return slide
    
    def create_data_insights_slide(self, insights_data: Any, source_refs: Dict[str, Any]) -> Any:
        """Create a branded data insights slide"""
        layout = self._get_layout_for_content('insights')
        slide = self.prs.slides.add_slide(layout)
        
        # Add title
        title_shape = self._add_branded_title(slide, "Key Insights")
        
        # Add insights with brand styling
        if insights_data:
            self._add_branded_insights_bullets(slide, insights_data, Inches(1), Inches(2))
        
        # Add source attribution
        self.add_source_attribution(slide, source_refs)
        
        return slide
    
    def _get_layout_for_content(self, content_type: str) -> Any:
        """Get appropriate slide layout for content type"""
        current_template = self.brand_manager.get_current_template()
        
        if current_template:
            # Try to find specific layout type
            if content_type == 'title':
                layout = current_template.get_layout_by_type('title_slide')
            else:
                layout = current_template.get_layout_by_type('content_slide')
            
            if layout:
                layout_index = layout.get('index', 0)
                if layout_index < len(self.prs.slide_layouts):
                    return self.prs.slide_layouts[layout_index]
        
        # Fallback to generic layouts
        if len(self.prs.slide_layouts) > 6:
            return self.prs.slide_layouts[6]  # Blank layout
        elif len(self.prs.slide_layouts) > 1:
            return self.prs.slide_layouts[1]  # Title and content
        else:
            return self.prs.slide_layouts[0]  # Title slide
    
    def _add_branded_title(self, slide: Any, title_text: str) -> Any:
        """Add title with brand-appropriate styling"""
        # Try to use title placeholder if available
        title_shape = None
        
        for shape in slide.shapes:
            try:
                # Check if shape is a placeholder first
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    # Now safe to access placeholder_format
                    if shape.placeholder_format and 'TITLE' in str(shape.placeholder_format.type):
                        title_shape = shape
                        break
            except:
                # Skip shapes that can't be checked
                continue
        
        # If no title placeholder, create text box
        if title_shape is None:
            title_shape = slide.shapes.add_textbox(
                Inches(1), Inches(0.5), Inches(8), Inches(1)
            )
        
        # Set title text (no need to clear, setting text replaces content)
        title_frame = title_shape.text_frame
        title_frame.text = title_text
        
        # Apply brand styling
        title_paragraph = title_frame.paragraphs[0]
        self._apply_font_style(title_paragraph, 'heading', size='large')
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        return title_shape
    
    def _add_branded_metrics_table(self, slide: Any, metrics_data: Dict[str, Any], 
                                 left: float, top: float):
        """Add metrics table with brand styling and superscript citations"""
        if not metrics_data:
            return
        
        # Calculate table dimensions
        rows = min(len(metrics_data) + 1, 10)  # +1 for header
        cols = 3  # Metric, Value, Source
        
        # Add table
        table_shape = slide.shapes.add_table(rows, cols, left, top, Inches(8), Inches(4))
        table = table_shape.table
        
        # Set column widths
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(2) 
        table.columns[2].width = Inches(3)
        
        # Add headers with brand styling
        headers = ['Metric', 'Value', 'Source']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._style_branded_header_cell(cell)
        
        # Track citations for footnotes
        citation_map = {}
        citation_counter = 1
        
        # Add data rows
        row_idx = 1
        for metric_name, metric_info in metrics_data.items():
            if row_idx >= rows:
                break
            
            # Metric name
            metric_cell = table.cell(row_idx, 0)
            metric_cell.text = str(metric_name)
            self._apply_cell_font_style(metric_cell)
            
            # Value with superscript citation
            value = metric_info.get('value', 'N/A')
            formatted_value = self._format_financial_value(value)
            value_cell = table.cell(row_idx, 1)
            
            # Get confidence level if available
            confidence = 1.0
            data_point_id = metric_info.get('data_point_id')
            if self.source_tracker and data_point_id:
                data_point = self.source_tracker.data_points.get(data_point_id)
                if data_point:
                    confidence = data_point.confidence
            
            # Add value with superscript citation
            value_cell.text = ""
            paragraph = value_cell.text_frame.paragraphs[0]
            
            # Add clickable value if available
            if self.source_tracker and data_point_id:
                hyperlink_url = self.source_tracker.get_source_hyperlink(data_point_id, formatted_value)
                run = paragraph.add_run()
                run.text = formatted_value
                try:
                    run.hyperlink.address = hyperlink_url
                    run.font.color.rgb = self._hex_to_rgb(self._get_brand_color('accent1', '#0066CC'))
                    run.font.underline = True
                except:
                    pass
                run.font.size = Pt(12)
                run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            else:
                run = paragraph.add_run()
                run.text = formatted_value
                run.font.size = Pt(12)
                run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            
            # Add superscript citation
            citation_run = paragraph.add_run()
            citation_run.text = f"[{citation_counter}]"
            citation_run.font.size = Pt(8)
            citation_run.font.superscript = True
            citation_run.font.color.rgb = self._get_confidence_color(confidence)
            
            # Store citation info
            citation_map[citation_counter] = {
                'metric': metric_name,
                'confidence': confidence,
                'data_point_id': data_point_id
            }
            
            # Source with enhanced attribution and confidence
            source_cell = table.cell(row_idx, 2)
            if self.source_tracker and data_point_id:
                # Use enhanced source attribution with confidence
                source_text = self.source_tracker.get_source_attribution_text(data_point_id, 'minimal')
                formatted_source, confidence_text = self._format_source_with_confidence(source_text, confidence)
                
                # Clear cell and add formatted text
                source_cell.text = ""
                paragraph = source_cell.text_frame.paragraphs[0]
                
                # Add clickable source
                run = paragraph.add_run()
                run.text = formatted_source
                try:
                    hyperlink_url = self.source_tracker.get_source_hyperlink(data_point_id, formatted_source)
                    run.hyperlink.address = hyperlink_url
                    run.font.color.rgb = self._hex_to_rgb(self._get_brand_color('secondary', '#666666'))
                    run.font.italic = True
                except:
                    pass
                run.font.size = Pt(10)
                run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
                
                # Add confidence indicator
                conf_run = paragraph.add_run()
                conf_run.text = f" {confidence_text}"
                conf_run.font.size = Pt(10)
                conf_run.font.color.rgb = self._get_confidence_color(confidence)
                conf_run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            else:
                # Fallback to basic cell reference
                cell_ref = metric_info.get('cell', 'Unknown')
                source_text = f"Cell {cell_ref}"
                source_cell.text = source_text
                self._apply_cell_font_style(source_cell, size='small')
            
            citation_counter += 1
            row_idx += 1
        
        # Store citations for later use in footnotes
        self._last_citations = citation_map
    
    def _add_branded_company_info(self, slide: Any, company_data: Dict[str, Any],
                                left: float, top: float):
        """Add company information with brand styling"""
        info_shape = slide.shapes.add_textbox(left, top, Inches(8), Inches(4))
        info_frame = info_shape.text_frame
        
        # Build company info text
        info_lines = []
        
        if 'name' in company_data:
            info_lines.append(f"Company: {company_data['name']}")
        
        if 'industry' in company_data:
            info_lines.append(f"Industry: {company_data['industry']}")
        
        if 'description' in company_data:
            info_lines.append(f"Description: {company_data['description']}")
        
        # Set text
        info_frame.text = '\\n\\n'.join(info_lines)
        
        # Apply brand styling
        for paragraph in info_frame.paragraphs:
            self._apply_font_style(paragraph, 'body', size='medium')
    
    def _add_branded_insights_bullets(self, slide: Any, insights: Any,
                                    left: float, top: float):
        """Add insights bullets with brand styling"""
        bullets_shape = slide.shapes.add_textbox(left, top, Inches(8), Inches(4))
        bullets_frame = bullets_shape.text_frame
        
        if isinstance(insights, list):
            for i, insight in enumerate(insights):
                if i == 0:
                    bullets_frame.text = f"• {insight}"
                    paragraph = bullets_frame.paragraphs[0]
                else:
                    paragraph = bullets_frame.add_paragraph()
                    paragraph.text = f"• {insight}"
                
                self._apply_font_style(paragraph, 'body', size='medium')
                
        elif isinstance(insights, dict):
            first = True
            for key, value in insights.items():
                text = f"• {key}: {value}"
                if first:
                    bullets_frame.text = text
                    paragraph = bullets_frame.paragraphs[0]
                    first = False
                else:
                    paragraph = bullets_frame.add_paragraph()
                    paragraph.text = text
                
                self._apply_font_style(paragraph, 'body', size='medium')
    
    def _style_branded_header_cell(self, cell: Any):
        """Apply brand styling to table header cells"""
        # Set background color to primary brand color
        cell.fill.solid()
        primary_color = self._get_brand_color('primary')
        cell.fill.fore_color.rgb = self._hex_to_rgb(primary_color)
        
        # Style text
        for paragraph in cell.text_frame.paragraphs:
            self._apply_font_style(paragraph, 'heading', size='small')
            # Set text color to white for contrast
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
    
    def _apply_cell_font_style(self, cell: Any, size: str = 'medium'):
        """Apply brand font styling to table cells"""
        for paragraph in cell.text_frame.paragraphs:
            self._apply_font_style(paragraph, 'body', size=size)
    
    def _apply_font_style(self, paragraph: Any, font_type: str = 'body', size: str = 'medium'):
        """Apply brand font styling to paragraph"""
        font_config = self.brand_config.get('fonts', {}).get(font_type, {})
        
        for run in paragraph.runs:
            # Set font family
            font_family = font_config.get('family', 'Calibri')
            run.font.name = font_family
            
            # Set font size
            size_key = f'size_{size}'
            font_size = font_config.get(size_key, font_config.get('size_medium', 14))
            run.font.size = Pt(font_size)
            
            # Set bold
            if font_config.get('bold', False):
                run.font.bold = True
            
            # Set color to dark brand color
            if font_type == 'heading':
                color = self._get_brand_color('dark1', '#000000')
            else:
                color = self._get_brand_color('dark1', '#333333')
            
            run.font.color.rgb = self._hex_to_rgb(color)
    
    def _get_brand_color(self, color_name: str, default: str = '#4F81BD') -> str:
        """Get brand color by name"""
        return self.brand_config.get('theme_colors', {}).get(color_name, default)
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor object"""
        # Remove # if present
        hex_color = hex_color.lstrip('#')
        
        # Convert to RGB
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        except:
            # Fallback to blue
            return RGBColor(79, 129, 189)
    
    def _format_financial_value(self, value: Any) -> str:
        """Format financial values for display"""
        if isinstance(value, (int, float)):
            if abs(value) > 1000000:
                return f"${value/1000000:.1f}M"
            elif abs(value) > 1000:
                return f"${value/1000:.0f}K"
            else:
                return f"${value:,.0f}"
        else:
            return str(value)
    
    def add_source_attribution(self, slide: Any, source_refs: Dict[str, Any]):
        """Add enhanced source attribution bar with icons, metadata, and brand styling"""
        if not source_refs:
            return
        
        # Create attribution bar background
        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.8),
            Inches(10), Inches(0.7)
        )
        
        # Style the bar with subtle gradient
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray background
        bar_shape.line.fill.background()  # No border
        
        # Add attribution content
        attr_shape = slide.shapes.add_textbox(
            Inches(0.3), Inches(6.85), Inches(9.4), Inches(0.6)
        )
        attr_frame = attr_shape.text_frame
        attr_frame.margin_left = 0
        attr_frame.margin_right = 0
        attr_frame.margin_top = Inches(0.05)
        attr_frame.margin_bottom = 0
        
        # Build enhanced attribution with metadata
        if self.source_tracker and isinstance(source_refs, dict):
            # Collect statistics
            data_point_ids = []
            doc_types = {'excel': 0, 'pdf': 0, 'word': 0}
            total_confidence = 0
            
            for source, details in source_refs.items():
                if isinstance(details, dict) and 'data_point_id' in details:
                    dp_id = details['data_point_id']
                    data_point_ids.append(dp_id)
                    
                    # Get data point info
                    if dp_id in self.source_tracker.data_points:
                        dp = self.source_tracker.data_points[dp_id]
                        total_confidence += dp.confidence
                        doc_types[dp.primary_source.document_type] += 1
            
            # Build attribution text with icons
            paragraph = attr_frame.paragraphs[0]
            
            # Add document icons and counts
            if doc_types['excel'] > 0:
                run = paragraph.add_run()
                run.text = "📊 "
                run.font.size = Pt(10)
                run = paragraph.add_run()
                run.text = f"Excel ({doc_types['excel']}) "
                run.font.size = Pt(9)
                run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
                run.font.color.rgb = RGBColor(100, 100, 100)
            
            if doc_types['pdf'] > 0:
                if doc_types['excel'] > 0:
                    run = paragraph.add_run()
                    run.text = " | "
                    run.font.size = Pt(9)
                    run.font.color.rgb = RGBColor(180, 180, 180)
                run = paragraph.add_run()
                run.text = "📄 "
                run.font.size = Pt(10)
                run = paragraph.add_run()
                run.text = f"PDF ({doc_types['pdf']}) "
                run.font.size = Pt(9)
                run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
                run.font.color.rgb = RGBColor(100, 100, 100)
            
            if doc_types['word'] > 0:
                if doc_types['excel'] > 0 or doc_types['pdf'] > 0:
                    run = paragraph.add_run()
                    run.text = " | "
                    run.font.size = Pt(9)
                    run.font.color.rgb = RGBColor(180, 180, 180)
                run = paragraph.add_run()
                run.text = "📝 "
                run.font.size = Pt(10)
                run = paragraph.add_run()
                run.text = f"Word ({doc_types['word']}) "
                run.font.size = Pt(9)
                run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
                run.font.color.rgb = RGBColor(100, 100, 100)
            
            # Add separator
            run = paragraph.add_run()
            run.text = "  •  "
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(180, 180, 180)
            
            # Add extraction metadata
            run = paragraph.add_run()
            run.text = f"Extracted from {len(set(source_refs.keys()))} document{'s' if len(source_refs) > 1 else ''} | "
            run.font.size = Pt(9)
            run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            run.font.color.rgb = RGBColor(100, 100, 100)
            
            # Add data points count
            run = paragraph.add_run()
            run.text = f"{len(data_point_ids)} data point{'s' if len(data_point_ids) > 1 else ''} | "
            run.font.size = Pt(9)
            run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            run.font.color.rgb = RGBColor(100, 100, 100)
            
            # Add average confidence
            if data_point_ids:
                avg_confidence = total_confidence / len(data_point_ids)
                run = paragraph.add_run()
                run.text = f"{int(avg_confidence * 100)}% avg confidence"
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
                run.font.color.rgb = self._get_confidence_color(avg_confidence)
            
            # Add "View Source Details >" link on the right
            run = paragraph.add_run()
            run.text = "    View Source Details >"
            run.font.size = Pt(9)
            run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            run.font.color.rgb = self._hex_to_rgb(self._get_brand_color('accent1', '#0066CC'))
            run.font.underline = True
            
            # Add citation footnotes if available
            if hasattr(self, '_last_citations') and self._last_citations:
                self._add_citation_footnotes(slide, self._last_citations)
        else:
            # Fallback to simple attribution
            attr_text = self._build_fallback_attribution(source_refs)
            attr_frame.text = attr_text
            
            # Style attribution text
            for paragraph in attr_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(9)
                    run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
                    run.font.color.rgb = RGBColor(100, 100, 100)
    
    def create_source_summary_slide(self, source_refs: Dict[str, Any]) -> Any:
        """Create a comprehensive source summary slide with statistics and confidence visualization"""
        try:
            layout = self._get_layout_for_content('content')
            slide = self.prs.slides.add_slide(layout)
            
            # Add title
            title_shape = self._add_branded_title(slide, "Data Sources & Methodology")
            
            if not self.source_tracker:
                # Fallback if no source tracker
                info_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
                info_shape.text_frame.text = "Source tracking not available"
                return slide
            
            # Collect comprehensive statistics
            stats = self._collect_source_statistics(source_refs)
            
            # Left side: Document list and extraction stats
            try:
                self._add_document_summary(slide, stats, Inches(0.5), Inches(1.8))
            except Exception as e:
                print(f"Error adding document summary: {e}")
            
            # Right side: Confidence distribution chart
            try:
                self._add_confidence_chart(slide, stats, Inches(5.5), Inches(1.8))
            except Exception as e:
                import traceback
                print(f"Error adding confidence chart: {e}")
                print(f"Traceback: {traceback.format_exc()}")
            
            # Bottom: Processing methodology
            try:
                self._add_methodology_section(slide, stats, Inches(0.5), Inches(5.2))
            except Exception as e:
                print(f"Error adding methodology section: {e}")
            
            return slide
        except Exception as e:
            print(f"Error creating source summary slide: {e}")
            # Return a basic slide with error message
            try:
                layout = self._get_layout_for_content('content')
                slide = self.prs.slides.add_slide(layout)
                title_shape = self._add_branded_title(slide, "Data Sources & Methodology")
                
                error_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
                error_shape.text_frame.text = f"Error creating source summary: {str(e)}"
                return slide
            except:
                # If even that fails, return None
                return None
    
    def _collect_source_statistics(self, source_refs: Dict[str, Any]) -> Dict[str, Any]:
        """Collect comprehensive statistics about sources"""
        stats = {
            'documents': {},
            'total_data_points': 0,
            'confidence_distribution': {'high': 0, 'medium': 0, 'low': 0, 'very_low': 0},
            'doc_types': {'excel': 0, 'pdf': 0, 'word': 0},
            'extraction_methods': {},
            'data_point_ids': []
        }
        
        # Collect data point information
        for source, details in source_refs.items():
            if isinstance(details, dict) and 'data_point_id' in details:
                dp_id = details['data_point_id']
                stats['data_point_ids'].append(dp_id)
                
                if dp_id in self.source_tracker.data_points:
                    dp = self.source_tracker.data_points[dp_id]
                    stats['total_data_points'] += 1
                    
                    # Document info
                    doc_name = dp.primary_source.document_name
                    if doc_name not in stats['documents']:
                        stats['documents'][doc_name] = {
                            'type': dp.primary_source.document_type,
                            'data_points': 0,
                            'pages_or_sheets': set(),
                            'avg_confidence': 0,
                            'total_confidence': 0
                        }
                    
                    stats['documents'][doc_name]['data_points'] += 1
                    stats['documents'][doc_name]['total_confidence'] += dp.confidence
                    
                    if dp.primary_source.page_or_sheet:
                        stats['documents'][doc_name]['pages_or_sheets'].add(dp.primary_source.page_or_sheet)
                    
                    # Document type counts
                    stats['doc_types'][dp.primary_source.document_type] += 1
                    
                    # Confidence distribution
                    if dp.confidence >= 0.9:
                        stats['confidence_distribution']['high'] += 1
                    elif dp.confidence >= 0.7:
                        stats['confidence_distribution']['medium'] += 1
                    elif dp.confidence >= 0.5:
                        stats['confidence_distribution']['low'] += 1
                    else:
                        stats['confidence_distribution']['very_low'] += 1
                    
                    # Extraction method
                    method = dp.primary_source.extraction_method or 'Standard'
                    stats['extraction_methods'][method] = stats['extraction_methods'].get(method, 0) + 1
        
        # Calculate averages
        for doc_info in stats['documents'].values():
            if doc_info['data_points'] > 0:
                doc_info['avg_confidence'] = doc_info['total_confidence'] / doc_info['data_points']
        
        return stats
    
    def _add_document_summary(self, slide: Any, stats: Dict[str, Any], left: float, top: float):
        """Add document summary section to source slide"""
        # Create container
        summary_shape = slide.shapes.add_textbox(left, top, Inches(4.5), Inches(3))
        summary_frame = summary_shape.text_frame
        
        # Add section title
        p = summary_frame.paragraphs[0]
        p.text = "Processed Documents"
        self._apply_font_style(p, 'heading', size='medium')
        p.space_after = Pt(12)
        
        # Add document list
        for doc_name, doc_info in stats['documents'].items():
            p = summary_frame.add_paragraph()
            
            # Document icon
            icon = {'excel': '📊', 'pdf': '📄', 'word': '📝'}.get(doc_info['type'], '📄')
            run = p.add_run()
            run.text = f"{icon} "
            run.font.size = Pt(12)
            
            # Document name
            run = p.add_run()
            run.text = doc_name
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            
            p.space_after = Pt(4)
            
            # Document stats
            p = summary_frame.add_paragraph()
            p.left_indent = Inches(0.3)
            run = p.add_run()
            pages_sheets = len(doc_info['pages_or_sheets'])
            run.text = f"• {doc_info['data_points']} data points from {pages_sheets} {'sheets' if doc_info['type'] == 'excel' else 'pages'}"
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(100, 100, 100)
            run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            
            # Confidence indicator
            run = p.add_run()
            run.text = f" • {int(doc_info['avg_confidence'] * 100)}% avg confidence"
            run.font.size = Pt(10)
            run.font.color.rgb = self._get_confidence_color(doc_info['avg_confidence'])
            run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            
            p.space_after = Pt(8)
        
        # Add totals
        p = summary_frame.add_paragraph()
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = f"Total: {stats['total_data_points']} data points from {len(stats['documents'])} documents"
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
    
    def _add_confidence_chart(self, slide: Any, stats: Dict[str, Any], left: float, top: float):
        """Add confidence distribution chart"""
        try:
            # Create pie chart data
            conf_dist = stats['confidence_distribution']
            data = {
                'High (≥90%)': conf_dist['high'],
                'Medium (70-89%)': conf_dist['medium'],
                'Low (50-69%)': conf_dist['low'],
                'Very Low (<50%)': conf_dist['very_low']
            }
            
            # Filter out zero values
            data = {k: v for k, v in data.items() if v > 0}
            
            if data and hasattr(self, 'chart_generator'):
                # Generate pie chart - using only supported parameters
                chart_buffer = self.chart_generator.create_pie_chart(
                    data, 
                    title='Confidence Distribution',
                    show_percentages=True,
                    size=(4, 3.5)
                )
                
                if chart_buffer:
                    # Add chart to slide
                    slide.shapes.add_picture(chart_buffer, left, top, width=Inches(4))
            else:
                # Fallback: Add text-based confidence summary
                chart_shape = slide.shapes.add_textbox(left, top, Inches(4), Inches(3))
                chart_frame = chart_shape.text_frame
                
                p = chart_frame.paragraphs[0]
                p.text = "Confidence Distribution"
                self._apply_font_style(p, 'heading', size='small')
                p.space_after = Pt(8)
                
                total_points = sum(conf_dist.values())
                if total_points > 0:
                    for level, count in data.items():
                        if count > 0:
                            p = chart_frame.add_paragraph()
                            percentage = int((count / total_points) * 100)
                            p.text = f"• {level}: {count} ({percentage}%)"
                            self._apply_font_style(p, 'body', size='small')
        except Exception as e:
            print(f"Error creating confidence chart: {e}")
            # Add fallback text
            try:
                chart_shape = slide.shapes.add_textbox(left, top, Inches(4), Inches(1))
                chart_frame = chart_shape.text_frame
                chart_frame.text = "Confidence chart unavailable"
            except:
                pass
    
    def _add_methodology_section(self, slide: Any, stats: Dict[str, Any], left: float, top: float):
        """Add methodology section"""
        method_shape = slide.shapes.add_textbox(left, top, Inches(9), Inches(1.3))
        method_frame = method_shape.text_frame
        
        # Title
        p = method_frame.paragraphs[0]
        p.text = "Extraction Methodology"
        self._apply_font_style(p, 'heading', size='small')
        p.space_after = Pt(8)
        
        # Methods used
        p = method_frame.add_paragraph()
        methods_text = []
        for method, count in stats['extraction_methods'].items():
            methods_text.append(f"{method} ({count})")
        
        run = p.add_run()
        run.text = f"Methods used: {', '.join(methods_text)}"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(100, 100, 100)
        run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
        
        # Quality notes
        if stats['confidence_distribution']['low'] > 0 or stats['confidence_distribution']['very_low'] > 0:
            p = method_frame.add_paragraph()
            run = p.add_run()
            low_conf_count = stats['confidence_distribution']['low'] + stats['confidence_distribution']['very_low']
            run.text = f"⚠ {low_conf_count} data points have low confidence and may require manual verification"
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(255, 140, 0)
            run.font.italic = True
            run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
    
    def create_title_slide(self, title: str, subtitle: str = None) -> Any:
        """Create a branded title slide with modern visual effects"""
        layout = self._get_layout_for_content('title')
        slide = self.prs.slides.add_slide(layout)
        
        # Apply simple but effective visual effects
        enhance_slide_simply(slide, 'title', self.brand_config)
        
        # Add main title with enhanced styling
        title_shape = self._add_branded_title(slide, title)
        
        # Apply enhanced title styling
        if title_shape and hasattr(title_shape, 'text_frame'):
            self.simple_effects.enhance_title_text(title_shape.text_frame, 'dramatic')
        
        # Add subtitle if provided
        if subtitle:
            subtitle_shape = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(8), Inches(1)
            )
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.text = subtitle
            
            # Style subtitle with modern effects
            for paragraph in subtitle_frame.paragraphs:
                self._apply_font_style(paragraph, 'body', size='large')
                paragraph.alignment = PP_ALIGN.CENTER
            
            # Apply enhanced subtitle styling
            self.simple_effects.enhance_title_text(subtitle_frame, 'subtitle')
        
        return slide
    
    def save_presentation(self, output_path: str) -> str:
        """Save the branded presentation"""
        self.prs.save(output_path)
        return output_path
    
    def get_available_templates(self) -> List[str]:
        """Get list of available brand templates"""
        return self.brand_manager.list_templates()
    
    def switch_template(self, template_name: str):
        """Switch to a different brand template"""
        self.brand_manager.set_current_template(template_name)
        self.brand_config = self.brand_manager.get_current_brand_config()
        # Reinitialize chart generator with adapted brand config
        chart_config = self._adapt_brand_config_for_charts(self.brand_config)
        self.chart_generator = ChartGenerator(chart_config)
        # Reinitialize presentation with new template
        self._init_presentation()
    
    def _add_clickable_value(self, cell: Any, display_text: str, data_point_id: str):
        """Add a clickable hyperlink to a table cell for a data value"""
        if not self.source_tracker:
            cell.text = display_text
            self._apply_cell_font_style(cell)
            return
        
        # Get hyperlink URL from source tracker
        hyperlink_url = self.source_tracker.get_source_hyperlink(data_point_id, display_text)
        
        # Clear cell and add hyperlink
        cell.text = ""
        paragraph = cell.text_frame.paragraphs[0]
        
        # Add hyperlink run
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.text = display_text
        
        try:
            # Add hyperlink to the run
            run.hyperlink.address = hyperlink_url
            
            # Style the hyperlink
            run.font.size = Pt(12)
            run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            run.font.color.rgb = self._hex_to_rgb(self._get_brand_color('accent1', '#0066CC'))  # Blue for links
            run.font.underline = True
            
        except Exception:
            # Fallback if hyperlink creation fails
            run.text = display_text
            self._apply_cell_font_style(cell)
    
    def _add_clickable_source(self, cell: Any, source_text: str, data_point_id: str):
        """Add a clickable source reference to a table cell"""
        if not self.source_tracker:
            cell.text = source_text
            self._apply_cell_font_style(cell, size='small')
            return
        
        # Get source context for tooltip/title
        source_context = self.source_tracker.get_source_context(data_point_id)
        tooltip_text = source_context.get('source_details', {}).get('location', 'Source details')
        
        # Get hyperlink URL
        hyperlink_url = self.source_tracker.get_source_hyperlink(data_point_id, source_text)
        
        # Clear cell and add hyperlink
        cell.text = ""
        paragraph = cell.text_frame.paragraphs[0]
        
        # Add hyperlink run
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.text = source_text
        
        try:
            # Add hyperlink
            run.hyperlink.address = hyperlink_url
            
            # Style the source link
            run.font.size = Pt(10)
            run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            run.font.color.rgb = self._hex_to_rgb(self._get_brand_color('secondary', '#666666'))
            run.font.italic = True
            
        except Exception:
            # Fallback if hyperlink creation fails
            run.text = source_text
            self._apply_cell_font_style(cell, size='small')
    
    def _build_fallback_attribution(self, source_refs: Any) -> str:
        """Build attribution text using fallback method for non-tracked sources"""
        if isinstance(source_refs, dict):
            attr_lines = []
            for source, details in source_refs.items():
                if isinstance(details, dict) and 'filename' in details:
                    attr_lines.append(f"Source: {details['filename']}")
                else:
                    attr_lines.append(f"Source: {source}")
            return " | ".join(attr_lines)
        elif isinstance(source_refs, list):
            return "Sources: " + ", ".join(str(ref) for ref in source_refs)
        else:
            return f"Source: {source_refs}"
    
    def _get_confidence_color(self, confidence: float) -> RGBColor:
        """Get color based on confidence level"""
        if confidence >= 0.9:
            return RGBColor(0, 128, 0)  # Green for high confidence
        elif confidence >= 0.7:
            return RGBColor(255, 165, 0)  # Orange for medium confidence
        elif confidence >= 0.5:
            return RGBColor(255, 140, 0)  # Dark orange for low confidence
        else:
            return RGBColor(255, 0, 0)  # Red for very low confidence
    
    def _add_superscript_citation(self, paragraph: Any, text: str, citation_num: int, 
                                confidence: Optional[float] = None):
        """Add text with superscript citation number"""
        # Add main text
        run = paragraph.add_run()
        run.text = text
        
        # Add superscript citation
        citation_run = paragraph.add_run()
        citation_run.text = f"[{citation_num}]"
        citation_run.font.size = Pt(8)
        citation_run.font.superscript = True
        
        # Color-code based on confidence if provided
        if confidence is not None:
            citation_run.font.color.rgb = self._get_confidence_color(confidence)
        else:
            citation_run.font.color.rgb = self._hex_to_rgb(self._get_brand_color('secondary', '#666666'))
    
    def _format_source_with_confidence(self, source_text: str, confidence: float) -> Tuple[str, str]:
        """Format source text with confidence indicator"""
        # Add confidence checkmarks
        if confidence >= 0.9:
            confidence_indicator = "✓✓✓"
        elif confidence >= 0.7:
            confidence_indicator = "✓✓"
        elif confidence >= 0.5:
            confidence_indicator = "✓"
        else:
            confidence_indicator = "⚠"
        
        # Format confidence percentage
        confidence_pct = f"{int(confidence * 100)}%"
        
        return source_text, f"{confidence_indicator} {confidence_pct}"
    
    def _add_citation_footnotes(self, slide: Any, citations: Dict[int, Dict[str, Any]]):
        """Add citation footnotes at the bottom of the slide"""
        if not citations:
            return
        
        # Add footnotes text box
        footnote_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.3), Inches(9), Inches(0.5)
        )
        footnote_frame = footnote_shape.text_frame
        footnote_frame.margin_top = 0
        footnote_frame.margin_bottom = 0
        
        paragraph = footnote_frame.paragraphs[0]
        
        # Add each citation
        for citation_num, citation_info in citations.items():
            if citation_num > 1:
                run = paragraph.add_run()
                run.text = "  "
            
            # Citation number
            run = paragraph.add_run()
            run.text = f"[{citation_num}]"
            run.font.size = Pt(7)
            run.font.color.rgb = self._get_confidence_color(citation_info['confidence'])
            run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
            
            # Source details
            if self.source_tracker and citation_info.get('data_point_id'):
                source_text = self.source_tracker.get_source_attribution_text(
                    citation_info['data_point_id'], 'minimal'
                )
                run = paragraph.add_run()
                run.text = f" {source_text}"
                run.font.size = Pt(7)
                run.font.color.rgb = RGBColor(120, 120, 120)
                run.font.italic = True
                run.font.name = self.brand_config.get('fonts', {}).get('body', {}).get('family', 'Calibri')
    
    def add_source_confidence_indicator(self, slide: Any, data_point_ids: List[str]):
        """Add visual indicators for source confidence levels"""
        if not self.source_tracker or not data_point_ids:
            return
        
        # Validate data consistency across multiple data points
        validation_results = self.source_tracker.validate_data_consistency(data_point_ids)
        
        if not validation_results['consistent']:
            # Add warning indicator for low confidence or inconsistent data
            warning_shape = slide.shapes.add_textbox(
                Inches(8.5), Inches(0.2), Inches(1.3), Inches(0.5)
            )
            warning_frame = warning_shape.text_frame
            warning_frame.text = "⚠ Review Sources"
            
            # Style warning
            for paragraph in warning_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(255, 140, 0)  # Orange warning
                    run.font.bold = True
    
    def create_chart_slide(self, title: str, chart_type: str, data: Any,
                          chart_options: Dict[str, Any] = None,
                          source_refs: Optional[Dict[str, Any]] = None) -> Any:
        """Create a slide with a chart
        
        Args:
            title: Slide title
            chart_type: Type of chart ('bar', 'line', 'pie', 'waterfall', 'scatter')
            data: Chart data (format depends on chart type)
            chart_options: Additional options for chart generation
            source_refs: Source attribution information
            
        Returns:
            Created slide object
        """
        # Get appropriate layout
        layout = self._get_layout_for_content('chart')
        slide = self.prs.slides.add_slide(layout)
        
        # Apply simple visual effects for chart slides
        enhance_slide_simply(slide, 'chart', self.brand_config)
        
        # Add title
        title_shape = self._add_branded_title(slide, title)
        
        # Generate chart based on type
        chart_options = chart_options or {}
        chart_buffer = None
        
        if chart_type == 'bar':
            chart_buffer = self.chart_generator.create_bar_chart(
                data, **chart_options
            )
        elif chart_type == 'line':
            chart_buffer = self.chart_generator.create_line_chart(
                data, **chart_options
            )
        elif chart_type == 'pie':
            chart_buffer = self.chart_generator.create_pie_chart(
                data, **chart_options
            )
        elif chart_type == 'waterfall':
            chart_buffer = self.chart_generator.create_waterfall_chart(
                data, **chart_options
            )
        elif chart_type == 'scatter':
            x_data = data.get('x', [])
            y_data = data.get('y', [])
            chart_buffer = self.chart_generator.create_scatter_plot(
                x_data, y_data, **chart_options
            )
        
        # Add chart image to slide
        if chart_buffer:
            left = Inches(1)
            top = Inches(2)
            height = Inches(4.5)
            
            pic = slide.shapes.add_picture(chart_buffer, left, top, height=height)
            
            # Center the image horizontally
            slide_width = self.prs.slide_width
            pic.left = int((slide_width - pic.width) / 2)
        
        # Add source attribution
        if source_refs:
            self.add_source_attribution(slide, source_refs)
        
        return slide
    
    def create_financial_dashboard_slide(self, financial_data: Dict[str, Any],
                                       source_refs: Optional[Dict[str, Any]] = None) -> Any:
        """Create a financial dashboard slide with multiple charts
        
        Args:
            financial_data: Dictionary containing different financial metrics
            source_refs: Source attribution information
            
        Returns:
            Created slide object
        """
        # Get blank layout for custom positioning
        layout = self._get_layout_for_content('blank')
        slide = self.prs.slides.add_slide(layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_shape.text_frame
        title_frame.text = "Financial Dashboard"
        title_paragraph = title_frame.paragraphs[0]
        self._apply_font_style(title_paragraph, 'heading', size='large')
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Create revenue chart (top left)
        if 'revenue' in financial_data:
            revenue_chart = self.chart_generator.create_bar_chart(
                financial_data['revenue'],
                title="Revenue by Quarter",
                size=(4, 3)
            )
            slide.shapes.add_picture(revenue_chart, Inches(0.5), Inches(1.5), height=Inches(2.5))
        
        # Create profit margin chart (top right)
        if 'profit_margin' in financial_data:
            margin_chart = self.chart_generator.create_line_chart(
                {'Profit Margin %': list(financial_data['profit_margin'].values())},
                title="Profit Margin Trend",
                x_values=list(financial_data['profit_margin'].keys()),
                size=(4, 3)
            )
            slide.shapes.add_picture(margin_chart, Inches(5), Inches(1.5), height=Inches(2.5))
        
        # Create expense breakdown (bottom left)
        if 'expenses' in financial_data:
            expense_chart = self.chart_generator.create_pie_chart(
                financial_data['expenses'],
                title="Expense Breakdown",
                show_percentages=True,
                size=(4, 3)
            )
            slide.shapes.add_picture(expense_chart, Inches(0.5), Inches(4.5), height=Inches(2.5))
        
        # Create cash flow waterfall (bottom right)
        if 'cash_flow' in financial_data:
            # Convert to list of tuples for waterfall chart
            cash_flow_data = [(k, v) for k, v in financial_data['cash_flow'].items()]
            waterfall_chart = self.chart_generator.create_waterfall_chart(
                cash_flow_data,
                title="Cash Flow Analysis",
                size=(4, 3)
            )
            slide.shapes.add_picture(waterfall_chart, Inches(5), Inches(4.5), height=Inches(2.5))
        
        # Add source attribution
        if source_refs:
            self.add_source_attribution(slide, source_refs)
        
        return slide
    
    def create_comparison_chart_slide(self, title: str, 
                                    comparison_data: Dict[str, Dict[str, float]],
                                    chart_type: str = 'bar',
                                    source_refs: Optional[Dict[str, Any]] = None) -> Any:
        """Create a slide comparing multiple data series
        
        Args:
            title: Slide title
            comparison_data: Dictionary of series names to data dictionaries
            chart_type: 'bar' or 'line'
            source_refs: Source attribution information
            
        Returns:
            Created slide object
        """
        layout = self._get_layout_for_content('chart')
        slide = self.prs.slides.add_slide(layout)
        
        # Add title
        title_shape = self._add_branded_title(slide, title)
        
        # Prepare data for grouped bar chart or multi-line chart
        if chart_type == 'bar':
            # For grouped bar chart, we need to restructure the data
            # This is a simplified version - in production, you'd use more sophisticated charting
            first_series = next(iter(comparison_data.values()))
            combined_data = {}
            
            for category in first_series.keys():
                values = []
                for series_name in comparison_data.keys():
                    if category in comparison_data[series_name]:
                        values.append(comparison_data[series_name][category])
                combined_data[f"{category}"] = sum(values) / len(values)  # Average for now
            
            chart_buffer = self.chart_generator.create_bar_chart(
                combined_data,
                title="",
                orientation="vertical",
                size=(8, 5)
            )
        else:  # line chart
            chart_buffer = self.chart_generator.create_line_chart(
                comparison_data,
                title="",
                size=(8, 5)
            )
        
        # Add chart to slide
        if chart_buffer:
            left = Inches(1)
            top = Inches(1.8)
            height = Inches(4.7)
            
            pic = slide.shapes.add_picture(chart_buffer, left, top, height=height)
            
            # Center the image
            slide_width = self.prs.slide_width
            pic.left = int((slide_width - pic.width) / 2)
        
        # Add source attribution
        if source_refs:
            self.add_source_attribution(slide, source_refs)
        
        return slide
    
    def _adapt_brand_config_for_charts(self, brand_config: Dict[str, Any]) -> Dict[str, Any]:
        """Adapt brand config structure for ChartGenerator compatibility"""
        adapted = brand_config.copy()
        
        # Handle fonts structure mismatch
        fonts = brand_config.get('fonts', {})
        if 'body' in fonts and isinstance(fonts['body'], dict):
            # Extract size from nested structure
            body_config = fonts['body']
            adapted['fonts'] = {
                'title_font': 'DejaVu Sans',  # Linux-compatible font
                'body_font': 'DejaVu Sans',   # Linux-compatible font
                'title_size': fonts.get('heading', {}).get('size_large', 16),
                'body_size': body_config.get('size_medium', 12)
            }
        
        # Handle colors structure - ChartGenerator expects 'colors' not 'theme_colors'
        if 'theme_colors' in brand_config and 'colors' not in brand_config:
            adapted['colors'] = brand_config['theme_colors']
        
        return adapted
    
    def create_thank_you_slide(self) -> Any:
        """Create a branded thank you/questions slide"""
        layout = self._get_layout_for_content('thank_you')
        slide = self.prs.slides.add_slide(layout)
        
        # Add background with brand colors
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        shape.fill.gradient()
        shape.fill.gradient_angle = 90
        
        # Use brand colors for gradient
        primary_color = self._get_brand_color('primary', '#4F81BD')
        bg_color = self._get_brand_color('background', '#FFFFFF')
        
        # Convert to RGB for gradient
        primary_rgb = self._hex_to_rgb(primary_color)
        bg_rgb = self._hex_to_rgb(bg_color)
        
        # Create darker version of primary color for gradient start
        shape.fill.gradient_stops[0].color.rgb = RGBColor(
            max(0, int(primary_color[1:3], 16) - 40),
            max(0, int(primary_color[3:5], 16) - 40), 
            max(0, int(primary_color[5:7], 16) - 40)
        )
        shape.fill.gradient_stops[1].color.rgb = bg_rgb
        shape.line.fill.background()
        
        # Add "Thank You" text with brand styling
        thank_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        thank_frame = thank_shape.text_frame
        thank_frame.text = "Thank You"
        
        p = thank_frame.paragraphs[0]
        self._apply_font_style(p, 'heading', size='extra_large')
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = self._hex_to_rgb(self._get_brand_color('primary', '#37407D'))
        
        # Add "Questions?" text
        questions_shape = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        questions_frame = questions_shape.text_frame
        questions_frame.text = "Questions?"
        
        p = questions_frame.paragraphs[0]
        self._apply_font_style(p, 'heading', size='large')
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = self._hex_to_rgb(self._get_brand_color('accent1', '#4F81BD'))
        
        return slide
    
    def add_slide_numbers(self, start_from=1, exclude_first=True, exclude_last=True):
        """Add slide numbers to all slides with brand styling"""
        total_slides = len(self.prs.slides)
        
        for idx, slide in enumerate(self.prs.slides):
            # Skip first slide (title) if requested
            if exclude_first and idx == 0:
                continue
            
            # Skip last slide (thank you) if requested
            if exclude_last and idx == total_slides - 1:
                continue
            
            # Add slide number with brand styling
            slide_num = idx if not exclude_first else idx
            footer_shape = slide.shapes.add_textbox(
                Inches(8.5), Inches(7), Inches(1), Inches(0.3)
            )
            footer_frame = footer_shape.text_frame
            footer_frame.text = f"{slide_num}/{total_slides - (1 if exclude_first else 0) - (1 if exclude_last else 0)}"
            
            footer_paragraph = footer_frame.paragraphs[0]
            footer_paragraph.font.size = Pt(10)
            footer_paragraph.alignment = PP_ALIGN.RIGHT
            
            # Apply brand colors to slide numbers
            accent_color = self._get_brand_color('accent2', '#808080')
            footer_paragraph.font.color.rgb = self._hex_to_rgb(accent_color)
        
        return slide
    
    # Rich Layout Methods for Substantial Slides
    def create_executive_dashboard(self, metrics: Dict[str, Any], 
                                 title: str = "Executive Dashboard") -> Any:
        """Create an executive dashboard slide with KPI cards"""
        layout = self._get_layout_for_content('content')
        slide = self.prs.slides.add_slide(layout)
        
        # Apply visual effects
        enhance_slide_simply(slide, 'content', self.brand_config)
        
        # Create the dashboard
        self.rich_layouts.create_executive_dashboard_slide(slide, metrics, title)
        
        return slide
    
    def create_multi_chart_analysis(self, data: Dict[str, Any],
                                  title: str = "Financial Analysis") -> Any:
        """Create a slide with multiple charts for comprehensive analysis"""
        layout = self._get_layout_for_content('content')
        slide = self.prs.slides.add_slide(layout)
        
        # Apply visual effects
        enhance_slide_simply(slide, 'chart', self.brand_config)
        
        # Create the multi-chart layout
        self.rich_layouts.create_multi_chart_analysis_slide(slide, data, title)
        
        return slide
    
    def create_timeline_roadmap(self, milestones: List[Dict[str, Any]] = None,
                              title: str = "Strategic Roadmap") -> Any:
        """Create a timeline/roadmap visualization"""
        layout = self._get_layout_for_content('content')
        slide = self.prs.slides.add_slide(layout)
        
        # Apply visual effects
        enhance_slide_simply(slide, 'content', self.brand_config)
        
        # Create the timeline
        self.rich_layouts.create_timeline_roadmap_slide(slide, milestones, title)
        
        return slide
    
    def create_comparison_matrix(self, comparison_data: Dict[str, Any] = None,
                               title: str = "Competitive Analysis") -> Any:
        """Create a comparison matrix slide"""
        layout = self._get_layout_for_content('content')
        slide = self.prs.slides.add_slide(layout)
        
        # Apply visual effects
        enhance_slide_simply(slide, 'content', self.brand_config)
        
        # Create the comparison matrix
        self.rich_layouts.create_comparison_matrix_slide(slide, comparison_data, title)
        
        return slide
    
    def create_swot_analysis(self, swot_data: Dict[str, List[str]] = None,
                           title: str = "SWOT Analysis") -> Any:
        """Create a SWOT analysis slide with 2x2 matrix"""
        layout = self._get_layout_for_content('content')
        slide = self.prs.slides.add_slide(layout)
        
        # Apply visual effects
        enhance_slide_simply(slide, 'content', self.brand_config)
        
        # Create the SWOT matrix
        self.rich_layouts.create_swot_analysis_slide(slide, swot_data, title)
        
        return slide