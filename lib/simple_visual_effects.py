#!/usr/bin/env python3
"""
Simple Visual Effects System for PowerPoint Presentations

This module provides robust, simple visual enhancements that work reliably
with the existing presentation generation system.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from typing import Dict, List, Any, Optional, Tuple
import random

class SimpleVisualEffects:
    """Simple, robust visual effects for presentations"""
    
    # Modern color palettes
    COLOR_PALETTES = {
        'corporate_elite': {
            'primary': '#1e3c72',
            'secondary': '#2a5298', 
            'accent': '#667eea',
            'light': '#f8f9fa'
        },
        'vibrant_tech': {
            'primary': '#4facfe',
            'secondary': '#00f2fe',
            'accent': '#43e97b',
            'light': '#f0f8ff'
        },
        'warm_minimal': {
            'primary': '#ffecd2',
            'secondary': '#fcb69f',
            'accent': '#d299c2',
            'light': '#fef9d7'
        }
    }
    
    def __init__(self, brand_config: Optional[Dict[str, Any]] = None):
        """Initialize with brand configuration"""
        self.brand_config = brand_config or {}
        self.current_palette = self._select_palette()
    
    def _select_palette(self) -> Dict[str, str]:
        """Select color palette based on brand config"""
        primary_color = self.brand_config.get('theme_colors', {}).get('primary', '#4F81BD')
        
        # Analyze primary color brightness to select appropriate palette
        rgb = self._hex_to_rgb(primary_color)
        brightness = sum(rgb) / 3
        
        if brightness > 180:
            return self.COLOR_PALETTES['warm_minimal']
        elif brightness > 100:
            return self.COLOR_PALETTES['corporate_elite']
        else:
            return self.COLOR_PALETTES['vibrant_tech']
    
    def enhance_slide_with_accent_bar(self, slide, position: str = 'right') -> None:
        """Add a simple accent bar to the slide"""
        try:
            accent_color = self.current_palette['accent']
            
            if position == 'right':
                # Right side accent bar
                accent_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(9.2), Inches(0.5),
                    Inches(0.3), Inches(6.5)
                )
            elif position == 'bottom':
                # Bottom accent bar
                accent_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.5), Inches(7.0),
                    Inches(9.0), Inches(0.2)
                )
            else:
                # Top accent bar (default)
                accent_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.5), Inches(0.3),
                    Inches(9.0), Inches(0.2)
                )
            
            # Apply solid color fill
            accent_shape.fill.solid()
            rgb = self._hex_to_rgb(accent_color)
            accent_shape.fill.fore_color.rgb = RGBColor(*rgb)
            
            # Remove border
            accent_shape.line.fill.background()
            
        except Exception as e:
            # If accent bar fails, continue silently
            pass
    
    def enhance_title_text(self, text_frame, style: str = 'dramatic') -> None:
        """Enhance title text with modern styling"""
        try:
            if not text_frame or not text_frame.paragraphs:
                return
                
            for paragraph in text_frame.paragraphs:
                if not paragraph.runs:
                    continue
                    
                for run in paragraph.runs:
                    try:
                        if style == 'dramatic':
                            run.font.size = Pt(36)
                            run.font.bold = True
                            # Use primary color
                            primary_color = self.current_palette['primary']
                            rgb = self._hex_to_rgb(primary_color)
                            run.font.color.rgb = RGBColor(*rgb)
                        elif style == 'subtitle':
                            run.font.size = Pt(16)
                            run.font.italic = True
                            # Use accent color
                            accent_color = self.current_palette['accent']
                            rgb = self._hex_to_rgb(accent_color)
                            run.font.color.rgb = RGBColor(*rgb)
                        
                        # Set modern font
                        run.font.name = 'Calibri Light'
                    except:
                        continue
                        
        except Exception as e:
            pass
    
    def add_subtle_background_shape(self, slide, shape_type: str = 'circle') -> None:
        """Add a subtle background decorative shape"""
        try:
            light_color = self.current_palette['light']
            
            if shape_type == 'circle':
                # Large subtle circle in background
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(7.0), Inches(-1.0),
                    Inches(5.0), Inches(5.0)
                )
            else:
                # Rounded rectangle
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(6.5), Inches(0.5),
                    Inches(3.0), Inches(6.0)
                )
            
            # Apply very light fill
            bg_shape.fill.solid()
            rgb = self._hex_to_rgb(light_color)
            bg_shape.fill.fore_color.rgb = RGBColor(*rgb)
            
            # Make it very transparent-looking
            # Note: PowerPoint transparency is tricky, so we use a very light color instead
            
            # Remove border
            bg_shape.line.fill.background()
            
            # Send to back
            try:
                # Move to beginning of shapes collection (send to back)
                shape_element = bg_shape.element
                shape_parent = shape_element.getparent()
                shape_parent.remove(shape_element)
                shape_parent.insert(0, shape_element)
            except:
                pass  # If reordering fails, shape will just be on top
                
        except Exception as e:
            pass
    
    def enhance_chart_area(self, slide) -> None:
        """Add subtle enhancement around chart area"""
        try:
            # Add subtle frame around typical chart area
            frame_color = self.current_palette['secondary']
            
            frame_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.0), Inches(1.5),
                Inches(8.0), Inches(5.0)
            )
            
            # Very light background
            frame_shape.fill.solid()
            rgb = self._hex_to_rgb(self.current_palette['light'])
            frame_shape.fill.fore_color.rgb = RGBColor(*rgb)
            
            # Subtle border
            frame_shape.line.color.rgb = self._hex_to_rgb(frame_color)
            frame_shape.line.width = Pt(1)
            
            # Send to back
            try:
                shape_element = frame_shape.element
                shape_parent = shape_element.getparent()
                shape_parent.remove(shape_element)
                shape_parent.insert(0, shape_element)
            except:
                pass
                
        except Exception as e:
            pass
    
    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        try:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except:
            return (79, 129, 189)  # Default blue
    
    def get_current_palette(self) -> Dict[str, str]:
        """Get the current color palette"""
        return self.current_palette.copy()

def enhance_slide_simply(slide, slide_type: str = 'content', brand_config: Optional[Dict[str, Any]] = None) -> None:
    """
    Simple function to enhance any slide with minimal visual effects
    
    Args:
        slide: PowerPoint slide object
        slide_type: Type of slide ('title', 'content', 'chart')
        brand_config: Brand configuration dictionary
    """
    effects = SimpleVisualEffects(brand_config)
    
    if slide_type == 'title':
        # Title slide enhancements
        effects.add_subtle_background_shape(slide, 'circle')
        effects.enhance_slide_with_accent_bar(slide, 'bottom')
    elif slide_type == 'chart':
        # Chart slide enhancements
        effects.enhance_chart_area(slide)
        effects.enhance_slide_with_accent_bar(slide, 'right')
    else:
        # Content slide enhancements
        effects.enhance_slide_with_accent_bar(slide, 'right')

# Example usage
if __name__ == "__main__":
    print("🎨 Simple Visual Effects System")
    print("Available color palettes:")
    
    effects = SimpleVisualEffects()
    for palette_name, colors in SimpleVisualEffects.COLOR_PALETTES.items():
        print(f"  - {palette_name}: {colors}")
    
    print("\n✨ Ready for simple, robust visual enhancements!")