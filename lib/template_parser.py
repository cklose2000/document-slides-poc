"""
PowerPoint Template Parser for Brand Consistency

This module extracts styling information from PowerPoint templates including:
- Theme colors, fonts, and layouts
- Master slide configurations  
- Layout dimensions and positioning
- Background styles and brand elements
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import json
from typing import Dict, List, Any, Optional, Tuple

class TemplateParser:
    """Parse PowerPoint templates to extract brand styling information"""
    
    def __init__(self, template_path: str):
        """Initialize parser with template file path"""
        self.template_path = template_path
        self.presentation = None
        self.brand_config = {}
        
        if os.path.exists(template_path):
            try:
                self.presentation = Presentation(template_path)
                self._parse_template()
            except Exception as e:
                raise ValueError(f"Failed to parse template: {str(e)}")
        else:
            raise FileNotFoundError(f"Template file not found: {template_path}")
    
    def _parse_template(self):
        """Parse the template and extract all brand information"""
        self.brand_config = {
            'template_path': self.template_path,
            'theme_colors': self._extract_theme_colors(),
            'fonts': self._extract_font_families(),
            'layouts': self._extract_layouts(),
            'master_slides': self._extract_master_slides(),
            'background_styles': self._extract_background_styles(),
            'slide_dimensions': self._get_slide_dimensions()
        }
    
    def _extract_theme_colors(self) -> Dict[str, str]:
        """Extract theme colors from the presentation"""
        theme_colors = {}
        
        try:
            # Access theme colors from the presentation's color scheme
            color_scheme = self.presentation.slide_master.color_scheme
            
            # PowerPoint has 12 standard theme colors
            color_names = [
                'dark1', 'light1', 'dark2', 'light2',
                'accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6',
                'hyperlink', 'followedHyperlink'
            ]
            
            for i, color_name in enumerate(color_names):
                try:
                    if hasattr(color_scheme, f'_{color_name}'):
                        color = getattr(color_scheme, f'_{color_name}')
                        if hasattr(color, 'rgb'):
                            rgb = color.rgb
                            theme_colors[color_name] = f"#{rgb:06x}"
                except:
                    continue
                    
        except Exception:
            # Fallback: extract colors from actual slides
            theme_colors = self._extract_colors_from_slides()
        
        # Ensure we have at least basic colors
        if not theme_colors:
            theme_colors = {
                'primary': '#4F81BD',      # Default blue
                'secondary': '#F79646',    # Default orange  
                'accent1': '#9BBB59',      # Default green
                'accent2': '#8064A2',      # Default purple
                'dark1': '#000000',        # Black
                'light1': '#FFFFFF'        # White
            }
            
        return theme_colors
    
    def _extract_colors_from_slides(self) -> Dict[str, str]:
        """Extract colors by analyzing actual slide content"""
        colors_found = {}
        color_frequency = {}
        
        try:
            # Look at first few slides to find common colors
            slide_count = min(len(self.presentation.slides), 3)
            for slide_idx in range(slide_count):
                try:
                    slide = self.presentation.slides[slide_idx]
                    for shape in slide.shapes:
                        try:
                            # Check fill colors
                            if hasattr(shape, 'fill') and shape.fill.type:
                                if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                                    rgb = shape.fill.fore_color.rgb
                                    color_hex = f"#{rgb:06x}"
                                    color_frequency[color_hex] = color_frequency.get(color_hex, 0) + 1
                            
                            # Check text colors
                            if hasattr(shape, 'text_frame'):
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'):
                                            rgb = run.font.color.rgb
                                            color_hex = f"#{rgb:06x}"
                                            color_frequency[color_hex] = color_frequency.get(color_hex, 0) + 1
                        except:
                            continue
                except:
                    continue
        except Exception:
            pass
        
        # Use most frequent colors as theme colors
        sorted_colors = sorted(color_frequency.items(), key=lambda x: x[1], reverse=True)
        
        if len(sorted_colors) >= 2:
            colors_found['primary'] = sorted_colors[0][0]
            colors_found['secondary'] = sorted_colors[1][0]
        
        if len(sorted_colors) >= 4:
            colors_found['accent1'] = sorted_colors[2][0]
            colors_found['accent2'] = sorted_colors[3][0]
        
        return colors_found
    
    def _extract_font_families(self) -> Dict[str, Dict[str, Any]]:
        """Extract font information from the template"""
        fonts = {
            'heading': {
                'family': 'Calibri',
                'size_large': 32,
                'size_medium': 24,
                'size_small': 18,
                'bold': True
            },
            'body': {
                'family': 'Calibri',
                'size_large': 16,
                'size_medium': 14,
                'size_small': 12,
                'bold': False
            }
        }
        
        try:
            # Check master slide for font themes
            master = self.presentation.slide_master
            
            # Look for text in master slide shapes
            for shape in master.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                # Determine if this is likely a heading or body font
                                if run.font.size and run.font.size.pt > 20:
                                    fonts['heading']['family'] = run.font.name
                                    fonts['heading']['size_large'] = int(run.font.size.pt)
                                else:
                                    fonts['body']['family'] = run.font.name
                                    fonts['body']['size_medium'] = int(run.font.size.pt) if run.font.size else 14
                                break
        except Exception:
            pass
        
        return fonts
    
    def _extract_layouts(self) -> List[Dict[str, Any]]:
        """Extract layout information for different slide types"""
        layouts = []
        
        for idx, layout in enumerate(self.presentation.slide_layouts):
            layout_info = {
                'index': idx,
                'name': layout.name if hasattr(layout, 'name') else f"Layout {idx}",
                'placeholders': [],
                'layout_type': self._classify_layout(layout)
            }
            
            # Extract placeholder information
            for placeholder in layout.placeholders:
                placeholder_info = {
                    'index': placeholder.placeholder_format.idx,
                    'type': str(placeholder.placeholder_format.type),
                    'left': placeholder.left,
                    'top': placeholder.top, 
                    'width': placeholder.width,
                    'height': placeholder.height
                }
                layout_info['placeholders'].append(placeholder_info)
            
            layouts.append(layout_info)
        
        return layouts
    
    def _classify_layout(self, layout) -> str:
        """Classify layout type based on placeholders"""
        placeholder_types = []
        
        for placeholder in layout.placeholders:
            placeholder_types.append(str(placeholder.placeholder_format.type))
        
        # Simple classification based on placeholder types
        if 'PP_PLACEHOLDER.TITLE' in placeholder_types and 'PP_PLACEHOLDER.SUBTITLE' in placeholder_types:
            return 'title_slide'
        elif 'PP_PLACEHOLDER.TITLE' in placeholder_types and 'PP_PLACEHOLDER.BODY' in placeholder_types:
            return 'content_slide'
        elif len(placeholder_types) == 0:
            return 'blank_slide'
        else:
            return 'custom_slide'
    
    def _extract_master_slides(self) -> Dict[str, Any]:
        """Extract master slide information"""
        master_info = {
            'background': None,
            'logo_position': None,
            'footer_info': None
        }
        
        try:
            master = self.presentation.slide_master
            
            # Check for background
            if hasattr(master, 'background'):
                master_info['background'] = 'custom_background'
            
            # Look for logo or consistent elements
            try:
                for shape in master.shapes:
                    if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        master_info['logo_position'] = {
                            'left': getattr(shape, 'left', 0),
                            'top': getattr(shape, 'top', 0),
                            'width': getattr(shape, 'width', 0),
                            'height': getattr(shape, 'height', 0)
                        }
                        break
            except Exception:
                pass
                    
        except Exception:
            pass
        
        return master_info
    
    def _extract_background_styles(self) -> Dict[str, Any]:
        """Extract background styling information"""
        background_info = {
            'type': 'solid_color',
            'color': '#FFFFFF',
            'has_watermark': False,
            'has_logo': False
        }
        
        try:
            # Check first slide for background pattern
            if self.presentation.slides:
                first_slide = self.presentation.slides[0]
                
                # Check for background fill
                if hasattr(first_slide.background, 'fill'):
                    fill = first_slide.background.fill
                    if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                        rgb = fill.fore_color.rgb
                        background_info['color'] = f"#{rgb:06x}"
                
                # Check for logo/watermark
                for shape in first_slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        background_info['has_logo'] = True
                        break
                        
        except Exception:
            pass
        
        return background_info
    
    def _get_slide_dimensions(self) -> Dict[str, float]:
        """Get slide dimensions in inches"""
        try:
            return {
                'width': self.presentation.slide_width / 914400,  # Convert EMU to inches
                'height': self.presentation.slide_height / 914400
            }
        except:
            return {
                'width': 10.0,  # Standard 4:3 ratio
                'height': 7.5
            }
    
    def get_brand_config(self) -> Dict[str, Any]:
        """Get the complete brand configuration"""
        return self.brand_config
    
    def get_theme_color(self, color_name: str, default: str = '#4F81BD') -> str:
        """Get a specific theme color by name"""
        return self.brand_config.get('theme_colors', {}).get(color_name, default)
    
    def get_font_config(self, font_type: str = 'body') -> Dict[str, Any]:
        """Get font configuration for heading or body text"""
        return self.brand_config.get('fonts', {}).get(font_type, {
            'family': 'Calibri',
            'size_medium': 14,
            'bold': False
        })
    
    def get_layout_by_type(self, layout_type: str) -> Optional[Dict[str, Any]]:
        """Get layout configuration by type"""
        layouts = self.brand_config.get('layouts', [])
        
        for layout in layouts:
            if layout.get('layout_type') == layout_type:
                return layout
        
        return None
    
    def save_brand_config(self, output_path: str):
        """Save brand configuration to JSON file"""
        with open(output_path, 'w') as f:
            json.dump(self.brand_config, f, indent=2, default=str)
    
    def load_brand_config(self, config_path: str):
        """Load brand configuration from JSON file"""
        with open(config_path, 'r') as f:
            self.brand_config = json.load(f)


class BrandManager:
    """Manage multiple brand templates and configurations"""
    
    def __init__(self, templates_dir: str = 'templates'):
        """Initialize brand manager with templates directory"""
        self.templates_dir = templates_dir
        self.templates = {}
        self.current_template = None
        
        # Create templates directory if it doesn't exist
        os.makedirs(templates_dir, exist_ok=True)
        
        # Load existing templates
        self._load_templates()
    
    def _load_templates(self):
        """Load all templates from the templates directory"""
        if not os.path.exists(self.templates_dir):
            return
        
        for filename in os.listdir(self.templates_dir):
            if filename.endswith('.pptx'):
                template_path = os.path.join(self.templates_dir, filename)
                template_name = filename.replace('.pptx', '')
                
                try:
                    parser = TemplateParser(template_path)
                    self.templates[template_name] = parser
                    
                    # Set first template as current if none set
                    if self.current_template is None:
                        self.current_template = template_name
                        
                except Exception as e:
                    print(f"Warning: Failed to load template {filename}: {str(e)}")
    
    def add_template(self, template_path: str, template_name: str = None) -> str:
        """Add a new template to the manager"""
        if template_name is None:
            template_name = os.path.basename(template_path).replace('.pptx', '')
        
        # Copy template to templates directory
        dest_path = os.path.join(self.templates_dir, f"{template_name}.pptx")
        
        try:
            import shutil
            shutil.copy2(template_path, dest_path)
            
            # Parse the template
            parser = TemplateParser(dest_path)
            self.templates[template_name] = parser
            
            # Set as current if first template
            if self.current_template is None:
                self.current_template = template_name
            
            return template_name
            
        except Exception as e:
            raise ValueError(f"Failed to add template: {str(e)}")
    
    def set_current_template(self, template_name: str):
        """Set the current active template"""
        if template_name in self.templates:
            self.current_template = template_name
        else:
            raise ValueError(f"Template '{template_name}' not found")
    
    def get_current_template(self) -> Optional[TemplateParser]:
        """Get the current active template parser"""
        if self.current_template and self.current_template in self.templates:
            return self.templates[self.current_template]
        return None
    
    def list_templates(self) -> List[str]:
        """List all available template names"""
        return list(self.templates.keys())
    
    def get_current_brand_config(self) -> Dict[str, Any]:
        """Get brand configuration for current template"""
        template = self.get_current_template()
        if template:
            return template.get_brand_config()
        
        # Return default configuration if no template
        return {
            'theme_colors': {
                'primary': '#4F81BD',
                'secondary': '#F79646',
                'accent1': '#9BBB59',
                'dark1': '#000000',
                'light1': '#FFFFFF'
            },
            'fonts': {
                'heading': {'family': 'Calibri', 'size_large': 32, 'bold': True},
                'body': {'family': 'Calibri', 'size_medium': 14, 'bold': False}
            }
        }