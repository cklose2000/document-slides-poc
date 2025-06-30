from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import tempfile
from io import BytesIO

# Import the new branded slide generator, chart generator, and layout engine
try:
    from .slide_generator_branded import BrandedSlideGenerator
    from .template_parser import BrandManager
    from .chart_generator import ChartGenerator
    from .layout_engine import SmartLayoutEngine
    BRANDED_AVAILABLE = True
    CHARTS_AVAILABLE = True
    LAYOUT_ENGINE_AVAILABLE = True
except ImportError:
    try:
        from slide_generator_branded import BrandedSlideGenerator
        from template_parser import BrandManager
        from chart_generator import ChartGenerator
        from layout_engine import SmartLayoutEngine
        BRANDED_AVAILABLE = True
        CHARTS_AVAILABLE = True
        LAYOUT_ENGINE_AVAILABLE = True
    except ImportError:
        BRANDED_AVAILABLE = False
        try:
            from .chart_generator import ChartGenerator
            from .layout_engine import SmartLayoutEngine
            CHARTS_AVAILABLE = True
            LAYOUT_ENGINE_AVAILABLE = True
        except ImportError:
            try:
                from chart_generator import ChartGenerator
                from layout_engine import SmartLayoutEngine
                CHARTS_AVAILABLE = True
                LAYOUT_ENGINE_AVAILABLE = True
            except ImportError:
                CHARTS_AVAILABLE = False
                LAYOUT_ENGINE_AVAILABLE = False

class SlideGenerator:
    def __init__(self, template_path='templates/firm_template.pptx', use_branding=True, 
                 source_tracker=None):
        """Initialize with template or create blank presentation"""
        self.template_path = template_path
        self.use_branding = use_branding and BRANDED_AVAILABLE
        self.source_tracker = source_tracker
        
        # Initialize chart generator if available
        if CHARTS_AVAILABLE:
            self.chart_generator = ChartGenerator()
        else:
            self.chart_generator = None
        
        # Initialize layout engine if available
        if LAYOUT_ENGINE_AVAILABLE:
            self.layout_engine = SmartLayoutEngine()
        else:
            self.layout_engine = None
        
        # Try to use branded generator if available
        if self.use_branding:
            try:
                self.brand_manager = BrandManager()
                # Check if template exists and add it
                if os.path.exists(template_path):
                    template_name = os.path.basename(template_path).replace('.pptx', '')
                    if template_name not in self.brand_manager.list_templates():
                        self.brand_manager.add_template(template_path, template_name)
                    self.brand_manager.set_current_template(template_name)
                
                # Pass source tracker to branded generator
                self.branded_generator = BrandedSlideGenerator(self.brand_manager, source_tracker=self.source_tracker)
                return
            except Exception as e:
                print(f"Warning: Could not initialize branded generator: {e}")
                self.use_branding = False
        
        # Fallback to original implementation
        if os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            # Create blank presentation if template doesn't exist
            self.prs = Presentation()
            self._setup_default_layouts()
    
    def _setup_default_layouts(self):
        """Set up default slide layouts if no template exists"""
        # This creates a basic presentation structure
        # In a real implementation, you'd want to create a proper template
        pass
    
    def create_financial_summary_slide(self, data, source_refs):
        """
        Create a slide with:
        - Title: Company Financial Summary
        - Key metrics in a table
        - Growth chart (if applicable)
        - Small source attribution text at bottom
        """
        # Use branded generator if available
        if self.use_branding:
            return self.branded_generator.create_financial_summary_slide(data, source_refs)
        
        # Get layout recommendation if layout engine is available
        layout_rec = None
        if self.layout_engine and data:
            content_data = {'financial_metrics': data} if isinstance(data, dict) else data
            layout_rec = self.layout_engine.recommend_layout(content_data)
            print(f"SMART LAYOUT ENGINE: {layout_rec.layout_type}")
            print(f"Font sizes: {layout_rec.font_sizes}")
            print(f"Element positions: {layout_rec.element_positions}")
            print(f"Should split: {layout_rec.split_recommendation}")
            print(f"Reasoning: {layout_rec.reasoning}")
            print(f"Data keys: {list(data.keys()) if isinstance(data, dict) else 'Not a dict'}")
            print(f"Has table position: {'table' in layout_rec.element_positions}")
            print(f"Has chart position: {'chart' in layout_rec.element_positions}")
        
        # Create separate slides for table and chart
        slides_created = []
        
        # First slide: Metrics Table
        slide_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[-1]
        table_slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title for table slide
        title_font_size = layout_rec.font_sizes.get('title', 36) if layout_rec else 36
        title_shape = table_slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1.2))
        title_frame = title_shape.text_frame
        title_frame.text = "Key Financial Metrics"
        
        # Style the title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(title_font_size)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.color.rgb = RGBColor(37, 64, 97)  # Dark blue
        
        # Add metrics table with full slide space
        if layout_rec and layout_rec.filtered_metrics:
            self._add_metrics_table_full_slide(table_slide, layout_rec.filtered_metrics, layout_rec)
        elif data and isinstance(data, dict):
            self._add_metrics_table_full_slide(table_slide, data, layout_rec)
        
        # Add source attribution
        self.add_source_attribution(table_slide, source_refs)
        slides_created.append(table_slide)
        
        # Second slide: Chart (if we have data for it)
        if self.chart_generator and data and isinstance(data, dict):
            metrics_to_chart = layout_rec.filtered_metrics if (layout_rec and layout_rec.filtered_metrics) else data
            
            # Check if we have enough data for a chart
            chart_data = {}
            for metric_name, metric_info in metrics_to_chart.items():
                if isinstance(metric_info, dict) and 'value' in metric_info:
                    value = self._extract_numeric_value(metric_info['value'])
                    if value is not None:
                        try:
                            from .layout_engine import MetricsPrioritizer
                            clean_name = MetricsPrioritizer.format_metric_name(metric_name)
                        except:
                            clean_name = self._clean_metric_name(metric_name)
                        chart_data[clean_name] = value
            
            if len(chart_data) >= 2:
                chart_slide = self.prs.slides.add_slide(slide_layout)
                
                # Add title for chart slide
                title_shape = chart_slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1.2))
                title_frame = title_shape.text_frame
                title_frame.text = "Financial Performance Chart"
                
                # Style the title
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.font.size = Pt(title_font_size)
                title_paragraph.font.bold = True
                title_paragraph.alignment = PP_ALIGN.CENTER
                title_paragraph.font.color.rgb = RGBColor(37, 64, 97)  # Dark blue
                
                # Add chart with full slide space
                self._add_chart_full_slide(chart_slide, metrics_to_chart, layout_rec)
                
                # Add source attribution
                self.add_source_attribution(chart_slide, source_refs)
                slides_created.append(chart_slide)
        
        return slides_created
    
    def _add_metrics_table_full_slide(self, slide, metrics_data, layout_rec):
        """Add a metrics table using full slide space with enhanced styling"""
        if not metrics_data:
            return
        
        # Use full slide area for table (center it nicely)
        left = 1.0
        top = 1.8
        width = 8.0
        height = 4.5
        
        rows = min(len(metrics_data) + 1, 10)  # +1 for header, max 10 rows
        cols = 4  # Metric, Value, Trend, Source
        
        # Add table shape
        table_shape = slide.shapes.add_table(
            rows, cols, 
            Inches(left), Inches(top), 
            Inches(width), Inches(height)
        )
        table = table_shape.table
        
        # Set column widths proportionally
        table.columns[0].width = Inches(3.2)  # Metric name
        table.columns[1].width = Inches(2.3)  # Value
        table.columns[2].width = Inches(0.8)  # Trend
        table.columns[3].width = Inches(1.7)  # Source
        
        # Add headers with good font sizing
        headers = ['Key Metrics', 'Value', 'Trend', 'Source']
        header_font_size = 16
        
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._style_header_cell(cell, Pt(header_font_size))
        
        # Add data rows with good spacing
        row_idx = 1
        body_font_size = 14
        
        for metric_name, metric_info in metrics_data.items():
            if row_idx >= rows:
                break
            
            # Apply alternating row colors
            for col in range(cols):
                cell = table.cell(row_idx, col)
                if row_idx % 2 == 0:
                    # Light gray background for even rows
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(248, 248, 248)
            
            # Metric name (clean it up)
            try:
                from .layout_engine import MetricsPrioritizer
                clean_name = MetricsPrioritizer.format_metric_name(metric_name)
            except:
                clean_name = str(metric_name).replace('_', ' ').title()
            
            table.cell(row_idx, 0).text = clean_name
            
            # Value (format properly)
            value = metric_info.get('value', 'N/A')
            numeric_value = self._extract_numeric_value(value) if isinstance(value, str) else value
            
            if isinstance(value, str) and any(c in value.lower() for c in ['m', 'k', '$']):
                formatted_value = f"${value}" if not value.startswith('$') else value
            elif isinstance(value, (int, float)):
                if abs(value) > 1000000:
                    formatted_value = f"${value/1000000:.1f}M"
                elif abs(value) > 1000:
                    formatted_value = f"${value/1000:.0f}K"
                else:
                    formatted_value = f"${value:,.0f}"
            else:
                formatted_value = str(value)
            
            table.cell(row_idx, 1).text = formatted_value
            
            # Trend indicator (based on metric type or value)
            trend_cell = table.cell(row_idx, 2)
            trend_text = ""
            trend_color = RGBColor(128, 128, 128)  # Default gray
            
            # Determine trend based on metric name or value
            metric_lower = metric_name.lower()
            if any(growth_term in metric_lower for growth_term in ['growth', 'increase', 'yoy', 'revenue', 'profit']):
                if numeric_value and numeric_value > 0:
                    trend_text = "↑"
                    trend_color = RGBColor(0, 128, 0)  # Green
                elif numeric_value and numeric_value < 0:
                    trend_text = "↓"
                    trend_color = RGBColor(255, 0, 0)  # Red
                else:
                    trend_text = "→"
            elif any(decrease_term in metric_lower for decrease_term in ['churn', 'cost', 'expense']):
                if numeric_value and numeric_value < 10:  # Low is good for these metrics
                    trend_text = "↓"
                    trend_color = RGBColor(0, 128, 0)  # Green
                else:
                    trend_text = "↑"
                    trend_color = RGBColor(255, 165, 0)  # Orange
            else:
                trend_text = "→"
            
            trend_cell.text = trend_text
            
            # Source
            source_info = metric_info.get('source', {})
            if isinstance(source_info, dict):
                doc_name = source_info.get('document', 'Unknown')
                if doc_name != 'Unknown':
                    source_text = doc_name.replace('.pdf', '').replace('.xlsx', '').replace('_', ' ').title()
                else:
                    source_text = 'Financial Report'
            else:
                source_text = 'Financial Report'
            
            table.cell(row_idx, 3).text = source_text
            
            # Style data cells
            for col in range(cols):
                cell = table.cell(row_idx, col)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(body_font_size)
                        if col == 1:  # Value column - make it bold
                            run.font.bold = True
                            # Color based on value
                            if numeric_value and numeric_value > 0:
                                run.font.color.rgb = RGBColor(0, 102, 51)  # Dark green
                            elif numeric_value and numeric_value < 0:
                                run.font.color.rgb = RGBColor(204, 0, 0)  # Dark red
                        elif col == 2:  # Trend column
                            run.font.size = Pt(18)  # Larger font for arrows
                            run.font.color.rgb = trend_color
                            paragraph.alignment = PP_ALIGN.CENTER
            
            row_idx += 1
    
    def _add_chart_full_slide(self, slide, metrics_data, layout_rec):
        """Add a chart using full slide space"""
        if not self.chart_generator:
            return
        
        try:
            # Prepare chart data from financial metrics
            chart_data = {}
            for metric_name, metric_info in metrics_data.items():
                if isinstance(metric_info, dict) and 'value' in metric_info:
                    value = self._extract_numeric_value(metric_info['value'])
                    if value is not None:
                        # Use the same clean name formatter as the table
                        try:
                            from .layout_engine import MetricsPrioritizer
                            clean_name = MetricsPrioritizer.format_metric_name(metric_name)
                        except:
                            clean_name = self._clean_metric_name(metric_name)
                        chart_data[clean_name] = value
            
            # Only create chart if we have 2+ metrics
            if len(chart_data) < 2:
                return
            
            # Use full slide area for chart
            left = 0.5
            top = 1.5
            width = 9.0
            height = 5.5
            
            # Generate horizontal bar chart for better readability
            chart_buffer = self.chart_generator.create_bar_chart(
                chart_data,
                title="",  # No title for cleaner look
                x_label="Value ($)",
                y_label="",
                orientation="horizontal",  # Horizontal bars for better metric names
                size=(width, height)
            )
            
            # Save chart to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_file.write(chart_buffer.getvalue())
                chart_path = tmp_file.name
            
            # Add chart to slide
            slide.shapes.add_picture(
                chart_path, 
                Inches(left), Inches(top), 
                Inches(width), Inches(height)
            )
            
            # Clean up temporary file
            try:
                os.unlink(chart_path)
            except:
                pass
            
        except Exception as e:
            print(f"Error adding full slide chart: {str(e)}")
    
    def _add_metrics_table_positioned(self, slide, metrics_data, layout_rec, table_pos):
        """Add a metrics table using smart layout positioning"""
        if not metrics_data:
            return
        
        left, top, width, height = table_pos
        rows = min(len(metrics_data) + 1, 10)  # +1 for header, max 10 rows
        cols = 3  # Metric, Value, Source
        
        # Add table shape with layout positioning
        table_shape = slide.shapes.add_table(
            rows, cols, 
            Inches(left), Inches(top), 
            Inches(width), Inches(height)
        )
        table = table_shape.table
        
        # Set column widths proportionally
        col_width = Inches(width / 3)
        table.columns[0].width = col_width
        table.columns[1].width = col_width
        table.columns[2].width = col_width
        
        # Add headers with layout-aware font sizing
        headers = ['Key Metrics', 'Value', 'Source Document']
        header_font_size = layout_rec.font_sizes.get('body', 14) - 2
        
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._style_header_cell(cell, Pt(max(header_font_size, 8)))
        
        # Add data rows with layout-aware styling
        row_idx = 1
        body_font_size = layout_rec.font_sizes.get('body', 14)
        
        for metric_name, metric_info in metrics_data.items():
            if row_idx >= rows:
                break
            
            # Metric name (clean it up)
            # Import the formatter at the top of the method if needed
            try:
                from .layout_engine import MetricsPrioritizer
                clean_name = MetricsPrioritizer.format_metric_name(metric_name)
            except:
                clean_name = str(metric_name).replace('_', ' ').title()
            
            table.cell(row_idx, 0).text = clean_name
            
            # Value (format properly)
            value = metric_info.get('value', 'N/A')
            if isinstance(value, str) and any(c in value.lower() for c in ['m', 'k', '$']):
                formatted_value = f"${value}" if not value.startswith('$') else value
            elif isinstance(value, (int, float)):
                if abs(value) > 1000000:
                    formatted_value = f"${value/1000000:.1f}M"
                elif abs(value) > 1000:
                    formatted_value = f"${value/1000:.0f}K"
                else:
                    formatted_value = f"${value:,.0f}"
            else:
                formatted_value = str(value)
            
            table.cell(row_idx, 1).text = formatted_value
            
            # Source
            source_info = metric_info.get('source', {})
            if isinstance(source_info, dict):
                doc_name = source_info.get('document', 'Unknown')
                if doc_name != 'Unknown':
                    source_text = doc_name.replace('.pdf', '').replace('.xlsx', '').replace('_', ' ').title()
                else:
                    source_text = 'Financial Report'
            else:
                source_text = 'Financial Report'
            
            table.cell(row_idx, 2).text = source_text
            
            # Style data cells with layout-aware font sizes
            for col in range(3):
                cell = table.cell(row_idx, col)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(max(body_font_size, 8))
                        if col == 1:  # Value column - make it bold
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0, 102, 51)  # Dark green
            
            row_idx += 1
    
    def _try_add_chart_positioned(self, slide, data, layout_rec, chart_pos):
        """Try to add a chart using smart layout positioning"""
        if not self.chart_generator:
            return False
        
        try:
            # Prepare chart data from financial metrics
            chart_data = {}
            for metric_name, metric_info in data.items():
                if isinstance(metric_info, dict) and 'value' in metric_info:
                    value = self._extract_numeric_value(metric_info['value'])
                    if value is not None:
                        # Use the same clean name formatter as the table
                        try:
                            from .layout_engine import MetricsPrioritizer
                            clean_name = MetricsPrioritizer.format_metric_name(metric_name)
                        except:
                            clean_name = self._clean_metric_name(metric_name)
                        chart_data[clean_name] = value
            
            # Only create chart if we have 2+ metrics
            if len(chart_data) < 2:
                return False
            
            left, top, width, height = chart_pos
            
            # Generate horizontal bar chart for better readability
            chart_buffer = self.chart_generator.create_bar_chart(
                chart_data,
                title="",  # No title for cleaner look
                x_label="Value ($)",
                y_label="",
                orientation="horizontal",  # Horizontal bars for better metric names
                size=(width, height)
            )
            
            # Save chart to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_file.write(chart_buffer.getvalue())
                chart_path = tmp_file.name
            
            # Add chart to slide with layout positioning
            slide.shapes.add_picture(
                chart_path, 
                Inches(left), Inches(top), 
                Inches(width), Inches(height)
            )
            
            # Clean up temporary file
            try:
                os.unlink(chart_path)
            except:
                pass
            
            return True
            
        except Exception as e:
            print(f"Error adding positioned chart to slide: {str(e)}")
            return False
    
    def _add_metrics_table(self, slide, metrics_data, left, top, table_width=Inches(8)):
        """Add a table with key metrics"""
        # Calculate table dimensions based on data
        if not metrics_data:
            return
        
        rows = min(len(metrics_data) + 1, 10)  # +1 for header, max 10 rows
        cols = 3  # Metric, Value, Source
        
        # Add table shape
        table_shape = slide.shapes.add_table(rows, cols, left, top, table_width, Inches(4))
        table = table_shape.table
        
        # Set column widths
        table.columns[0].width = Inches(3)  # Metric name
        table.columns[1].width = Inches(2.5)  # Value
        table.columns[2].width = Inches(2.5)  # Source
        
        # Add headers
        headers = ['Key Metrics', 'Value', 'Source Document']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._style_header_cell(cell)
        
        # Add data rows
        row_idx = 1
        for metric_name, metric_info in metrics_data.items():
            if row_idx >= rows:
                break
            
            # Metric name (clean it up)
            # Import the formatter at the top of the method if needed
            try:
                from .layout_engine import MetricsPrioritizer
                clean_name = MetricsPrioritizer.format_metric_name(metric_name)
            except:
                clean_name = str(metric_name).replace('_', ' ').title()
            
            table.cell(row_idx, 0).text = clean_name
            
            # Value (format properly)
            value = metric_info.get('value', 'N/A')
            if isinstance(value, str) and any(c in value.lower() for c in ['m', 'k', '$']):
                # Already formatted
                formatted_value = f"${value}" if not value.startswith('$') else value
            elif isinstance(value, (int, float)):
                if abs(value) > 1000000:
                    formatted_value = f"${value/1000000:.1f}M"
                elif abs(value) > 1000:
                    formatted_value = f"${value/1000:.0f}K"
                else:
                    formatted_value = f"${value:,.0f}"
            else:
                formatted_value = str(value)
            
            table.cell(row_idx, 1).text = formatted_value
            
            # Source (get proper document name)
            source_info = metric_info.get('source', {})
            if isinstance(source_info, dict):
                doc_name = source_info.get('document', 'Unknown')
                if doc_name != 'Unknown':
                    # Clean up filename
                    source_text = doc_name.replace('.pdf', '').replace('.xlsx', '').replace('_', ' ').title()
                else:
                    source_text = 'Financial Report'
            else:
                source_text = 'Financial Report'
            
            table.cell(row_idx, 2).text = source_text
            
            # Style data cells
            for col in range(3):
                cell = table.cell(row_idx, col)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
                        if col == 1:  # Value column - make it bold
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0, 102, 51)  # Dark green
            
            row_idx += 1
    
    def _style_header_cell(self, cell, font_size=None):
        """Apply styling to header cells"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(79, 129, 189)  # Blue background
        
        # Text styling
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)  # White text
                run.font.bold = True
                run.font.size = font_size if font_size else Pt(12)
    
    def create_company_overview_slide(self, company_data, source_refs):
        """Create a company overview slide with enhanced visual design"""
        # Use branded generator if available
        if self.use_branding:
            return self.branded_generator.create_company_overview_slide(company_data, source_refs)
        
        # Fallback to original implementation
        slide_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[-1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title with better styling
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1.2))
        title_frame = title_shape.text_frame
        title_frame.text = "Company Overview"
        
        # Style the title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(36)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.color.rgb = RGBColor(37, 64, 97)  # Dark blue
        
        # Add a decorative shape for visual interest
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9)
        height = Inches(0.02)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(79, 129, 189)  # Blue accent
        shape.line.fill.background()
        
        # Add company info with enhanced layout
        if company_data:
            self._add_company_info_enhanced(slide, company_data, Inches(1), Inches(2.2))
        
        # Add source attribution
        self.add_source_attribution(slide, source_refs)
        
        return slide
    
    def _add_company_info_enhanced(self, slide, company_data, left, top):
        """Add company information with enhanced visual layout including callout boxes"""
        current_top = top
        
        # Company Name (Large)
        if 'name' in company_data:
            name_shape = slide.shapes.add_textbox(left, current_top, Inches(8), Inches(0.8))
            name_frame = name_shape.text_frame
            name_frame.text = company_data['name']
            name_paragraph = name_frame.paragraphs[0]
            name_paragraph.font.size = Pt(28)
            name_paragraph.font.bold = True
            name_paragraph.font.color.rgb = RGBColor(37, 64, 97)
            current_top += Inches(1.0)
        
        # Industry (Medium with background)
        if 'industry' in company_data:
            # Add background shape
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left - Inches(0.2), current_top - Inches(0.1),
                Inches(8.4), Inches(0.8)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
            bg_shape.line.fill.background()
            
            industry_shape = slide.shapes.add_textbox(left, current_top, Inches(8), Inches(0.6))
            industry_frame = industry_shape.text_frame
            industry_frame.text = f"Industry: {company_data['industry']}"
            industry_paragraph = industry_frame.paragraphs[0]
            industry_paragraph.font.size = Pt(20)
            industry_paragraph.font.color.rgb = RGBColor(79, 129, 189)
            current_top += Inches(1.2)
        
        # Description (Body text) - shortened to make room for callout boxes
        if 'description' in company_data:
            desc_shape = slide.shapes.add_textbox(left, current_top, Inches(8), Inches(1.5))
            desc_frame = desc_shape.text_frame
            desc_frame.word_wrap = True
            desc_frame.text = company_data['description'][:200] + "..." if len(company_data.get('description', '')) > 200 else company_data.get('description', '')
            desc_paragraph = desc_frame.paragraphs[0]
            desc_paragraph.font.size = Pt(16)
            desc_paragraph.line_spacing = 1.5
            desc_paragraph.font.color.rgb = RGBColor(64, 64, 64)
            current_top += Inches(1.8)
        
        # Add key stats in callout boxes
        self._add_company_stats_callouts(slide, company_data, left, current_top)
    
    def _add_company_stats_callouts(self, slide, company_data, left, top):
        """Add key company statistics in visual callout boxes"""
        # Sample stats that might be in company data or defaults
        stats = [
            ('Global Reach', '15+ Countries', RGBColor(79, 129, 189)),
            ('Team Size', '450+ People', RGBColor(146, 208, 80)),
            ('Founded', '2019', RGBColor(255, 192, 0))
        ]
        
        # Create callout boxes for stats
        box_width = 2.5
        box_height = 1.0
        spacing = 0.3
        current_left = left
        
        for i, (label, value, color) in enumerate(stats[:3]):  # Max 3 stats
            # Create rounded rectangle callout box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                current_left, top,
                Inches(box_width), Inches(box_height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.fill.transparency = 0.85  # Make it subtle
            box.line.color.rgb = color
            box.line.width = Pt(2)
            
            # Add label
            label_shape = slide.shapes.add_textbox(
                current_left + Inches(0.1), top + Inches(0.1),
                Inches(box_width - 0.2), Inches(0.3)
            )
            label_frame = label_shape.text_frame
            label_frame.text = label
            label_paragraph = label_frame.paragraphs[0]
            label_paragraph.font.size = Pt(12)
            label_paragraph.font.color.rgb = RGBColor(64, 64, 64)
            label_paragraph.alignment = PP_ALIGN.CENTER
            
            # Add value
            value_shape = slide.shapes.add_textbox(
                current_left + Inches(0.1), top + Inches(0.4),
                Inches(box_width - 0.2), Inches(0.4)
            )
            value_frame = value_shape.text_frame
            value_frame.text = value
            value_paragraph = value_frame.paragraphs[0]
            value_paragraph.font.size = Pt(22)
            value_paragraph.font.bold = True
            value_paragraph.font.color.rgb = color
            value_paragraph.alignment = PP_ALIGN.CENTER
            
            current_left += Inches(box_width + spacing)
    
    def _add_company_info(self, slide, company_data, left, top):
        """Add company information text box (legacy method)"""
        info_shape = slide.shapes.add_textbox(left, top, Inches(8), Inches(4))
        info_frame = info_shape.text_frame
        
        # Build company info text
        info_lines = []
        
        if 'name' in company_data:
            info_lines.append(f"Company: {company_data['name']}")
        
        if 'industry' in company_data:
            info_lines.append(f"Industry: {company_data['industry']}")
        
        if 'description' in company_data:
            info_lines.append(f"Description: {company_data['description']}")
        
        info_frame.text = '\n\n'.join(info_lines)
        
        # Style the text
        for paragraph in info_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.space_after = Pt(12)
    
    def add_source_attribution(self, slide, source_refs):
        """Add small text at bottom with source references"""
        if not source_refs:
            return
        
        # Add attribution text box at bottom of slide
        attr_shape = slide.shapes.add_textbox(
            Inches(0.5), 
            Inches(6.5), 
            Inches(9), 
            Inches(1)
        )
        attr_frame = attr_shape.text_frame
        
        # Build attribution text
        if isinstance(source_refs, dict):
            attr_lines = []
            for source, details in source_refs.items():
                if isinstance(details, dict) and 'filename' in details:
                    attr_lines.append(f"Source: {details['filename']}")
                else:
                    attr_lines.append(f"Source: {source}")
            
            attr_text = " | ".join(attr_lines)
        elif isinstance(source_refs, list):
            attr_text = "Sources: " + ", ".join(str(ref) for ref in source_refs)
        else:
            attr_text = f"Source: {source_refs}"
        
        attr_frame.text = attr_text
        
        # Style attribution text (small and gray)
        for paragraph in attr_frame.paragraphs:
            paragraph.font.size = Pt(8)
            paragraph.font.color.rgb = RGBColor(128, 128, 128)  # Gray
            paragraph.alignment = PP_ALIGN.CENTER
    
    def create_data_insights_slide(self, insights_data, source_refs):
        """Create a slide with key data insights with enhanced visual design"""
        # Use branded generator if available
        if self.use_branding:
            return self.branded_generator.create_data_insights_slide(insights_data, source_refs)
        
        # Fallback to original implementation
        slide_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[-1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title with better styling
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1.2))
        title_frame = title_shape.text_frame
        title_frame.text = "Key Business Insights"
        
        # Style the title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(36)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.color.rgb = RGBColor(37, 64, 97)  # Dark blue
        
        # Add a decorative shape for visual interest
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9)
        height = Inches(0.02)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(79, 129, 189)  # Blue accent
        shape.line.fill.background()
        
        # Add insights with enhanced visual design
        if insights_data:
            self._add_insights_enhanced(slide, insights_data, Inches(0.8), Inches(2.2))
        
        # Add source attribution
        self.add_source_attribution(slide, source_refs)
        
        return slide
    
    def _add_insights_enhanced(self, slide, insights, left, top):
        """Add insights with enhanced visual design in a 2x2 grid layout"""
        if not isinstance(insights, list) or not insights:
            return
        
        # Use 2x2 grid layout for up to 4 insights
        insights_to_show = insights[:4]  # Maximum 4 insights
        
        # Grid parameters
        box_width = 4.0
        box_height = 1.8
        h_spacing = 0.5
        v_spacing = 0.4
        
        # Icon colors for different insight types
        colors = [
            RGBColor(79, 129, 189),   # Blue
            RGBColor(146, 208, 80),   # Green
            RGBColor(255, 192, 0),    # Yellow
            RGBColor(255, 102, 102)   # Red
        ]
        
        for i, insight in enumerate(insights_to_show):
            # Calculate position in grid
            row = i // 2
            col = i % 2
            
            box_left = left + (box_width + h_spacing) * col
            box_top = top + (box_height + v_spacing) * row
            
            # Create insight box with shadow effect
            # Shadow first
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                box_left + Inches(0.05), box_top + Inches(0.05),
                Inches(box_width), Inches(box_height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(200, 200, 200)
            shadow.line.fill.background()
            
            # Main box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                box_left, box_top,
                Inches(box_width), Inches(box_height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
            box.line.color.rgb = colors[i % len(colors)]
            box.line.width = Pt(2)
            
            # Add icon placeholder (colored circle with number)
            icon_size = 0.6
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                box_left + Inches(0.2), box_top + Inches(0.2),
                Inches(icon_size), Inches(icon_size)
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = colors[i % len(colors)]
            icon.line.fill.background()
            
            # Add number in icon
            icon_text = slide.shapes.add_textbox(
                box_left + Inches(0.2), box_top + Inches(0.2),
                Inches(icon_size), Inches(icon_size)
            )
            icon_frame = icon_text.text_frame
            icon_frame.text = str(i + 1)
            icon_paragraph = icon_frame.paragraphs[0]
            icon_paragraph.font.size = Pt(20)
            icon_paragraph.font.bold = True
            icon_paragraph.font.color.rgb = RGBColor(255, 255, 255)
            icon_paragraph.alignment = PP_ALIGN.CENTER
            icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Add insight text
            text_shape = slide.shapes.add_textbox(
                box_left + Inches(1.0), box_top + Inches(0.3),
                Inches(box_width - 1.2), Inches(box_height - 0.6)
            )
            text_frame = text_shape.text_frame
            text_frame.word_wrap = True
            text_frame.text = insight[:150] + "..." if len(insight) > 150 else insight
            
            text_paragraph = text_frame.paragraphs[0]
            text_paragraph.font.size = Pt(14)
            text_paragraph.line_spacing = 1.2
            text_paragraph.font.color.rgb = RGBColor(64, 64, 64)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def _add_insights_linear(self, slide, insights, left, top):
        """Add insights in a linear layout (legacy method)"""
        current_top = top
        
        if isinstance(insights, list):
            for i, insight in enumerate(insights):
                # Add background shape for each insight
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left - Inches(0.2), current_top - Inches(0.1),
                    Inches(8.6), Inches(0.9)
                )
                bg_shape.fill.solid()
                # Alternate colors for visual interest
                if i % 2 == 0:
                    bg_shape.fill.fore_color.rgb = RGBColor(245, 250, 255)  # Very light blue
                else:
                    bg_shape.fill.fore_color.rgb = RGBColor(250, 250, 250)  # Very light gray
                bg_shape.line.fill.background()
                
                # Add bullet point shape
                bullet_shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    left, current_top + Inches(0.15),
                    Inches(0.3), Inches(0.3)
                )
                bullet_shape.fill.solid()
                bullet_shape.fill.fore_color.rgb = RGBColor(79, 129, 189)  # Blue
                bullet_shape.line.fill.background()
                
                # Add insight text
                text_shape = slide.shapes.add_textbox(
                    left + Inches(0.5), current_top,
                    Inches(7.5), Inches(0.7)
                )
                text_frame = text_shape.text_frame
                text_frame.text = insight
                text_frame.word_wrap = True
                
                paragraph = text_frame.paragraphs[0]
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(37, 64, 97)
                
                current_top += Inches(1.1)
        
        elif isinstance(insights, dict):
            for i, (key, value) in enumerate(insights.items()):
                # Similar layout but with key-value format
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left - Inches(0.2), current_top - Inches(0.1),
                    Inches(8.6), Inches(0.9)
                )
                bg_shape.fill.solid()
                if i % 2 == 0:
                    bg_shape.fill.fore_color.rgb = RGBColor(245, 250, 255)
                else:
                    bg_shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
                bg_shape.line.fill.background()
                
                # Add bullet point
                bullet_shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    left, current_top + Inches(0.15),
                    Inches(0.3), Inches(0.3)
                )
                bullet_shape.fill.solid()
                bullet_shape.fill.fore_color.rgb = RGBColor(79, 129, 189)
                bullet_shape.line.fill.background()
                
                # Add text
                text_shape = slide.shapes.add_textbox(
                    left + Inches(0.5), current_top,
                    Inches(7.5), Inches(0.7)
                )
                text_frame = text_shape.text_frame
                text_frame.text = f"{key}: {value}"
                text_frame.word_wrap = True
                
                paragraph = text_frame.paragraphs[0]
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(37, 64, 97)
                
                current_top += Inches(1.1)
    
    def _add_insights_bullets(self, slide, insights, left, top):
        """Add insights as bullet points with better styling"""
        bullets_shape = slide.shapes.add_textbox(left, top, Inches(8), Inches(4))
        bullets_frame = bullets_shape.text_frame
        bullets_frame.margin_left = Inches(0.2)
        bullets_frame.margin_top = Inches(0.1)
        
        if isinstance(insights, list):
            for i, insight in enumerate(insights):
                if i == 0:
                    bullets_frame.text = f"• {insight}"
                    # Style first paragraph
                    p = bullets_frame.paragraphs[0]
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.space_after = Pt(12)
                    p.font.color.rgb = RGBColor(37, 64, 97)  # Dark blue
                else:
                    p = bullets_frame.add_paragraph()
                    p.text = f"• {insight}"
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.space_before = Pt(8)
                    p.space_after = Pt(8)
                    p.font.color.rgb = RGBColor(37, 64, 97)  # Dark blue
        elif isinstance(insights, dict):
            first = True
            i = 0
            for key, value in insights.items():
                text = f"• {key}: {value}"
                if first:
                    bullets_frame.text = text
                    p = bullets_frame.paragraphs[0]
                    first = False
                else:
                    p = bullets_frame.add_paragraph()
                    p.text = text
                
                p.font.size = Pt(18)
                p.font.bold = True
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.font.color.rgb = RGBColor(37, 64, 97)  # Dark blue
                i += 1
    
    def _try_add_chart(self, slide, data, left, top, width, height):
        """Try to add a chart to the slide based on financial data"""
        if not self.chart_generator:
            return False
        
        try:
            # Prepare chart data from financial metrics
            chart_data = {}
            for metric_name, metric_info in data.items():
                if isinstance(metric_info, dict) and 'value' in metric_info:
                    value = self._extract_numeric_value(metric_info['value'])
                    if value is not None:
                        clean_name = self._clean_metric_name(metric_name)
                        chart_data[clean_name] = value
            
            # Only create chart if we have 2+ metrics
            if len(chart_data) < 2:
                return False
            
            # Generate chart
            chart_buffer = self.chart_generator.create_bar_chart(
                chart_data,
                title="Key Metrics",
                y_label="Value ($)",
                size=(width.inches, height.inches)
            )
            
            # Save chart to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_file.write(chart_buffer.getvalue())
                chart_path = tmp_file.name
            
            # Add chart to slide
            slide.shapes.add_picture(chart_path, left, top, width, height)
            
            # Clean up temporary file
            try:
                os.unlink(chart_path)
            except:
                pass
            
            return True
            
        except Exception as e:
            print(f"Error adding chart to slide: {str(e)}")
            return False
    
    def _extract_numeric_value(self, value_str):
        """Extract numeric value from string (handles $, M, K, etc.)"""
        if not value_str:
            return None
        
        import re
        value_str = str(value_str).strip()
        
        # Remove currency symbols and spaces
        cleaned = re.sub(r'[$,\s]', '', value_str)
        
        # Handle multipliers
        multiplier = 1
        if cleaned.lower().endswith('m'):
            multiplier = 1_000_000
            cleaned = cleaned[:-1]
        elif cleaned.lower().endswith('k'):
            multiplier = 1_000
            cleaned = cleaned[:-1]
        elif cleaned.lower().endswith('b'):
            multiplier = 1_000_000_000
            cleaned = cleaned[:-1]
        
        # Extract number
        number_match = re.search(r'([\d.]+)', cleaned)
        if number_match:
            try:
                return float(number_match.group(1)) * multiplier
            except ValueError:
                return None
        
        return None
    
    def _clean_metric_name(self, name):
        """Clean and format metric names for display"""
        # Remove underscores and capitalize
        cleaned = name.replace('_', ' ').title()
        
        # Handle common abbreviations
        replacements = {
            'Yoy': 'YoY',
            'Roi': 'ROI',
            'Ebitda': 'EBITDA',
            'Arr': 'ARR',
            'Mrr': 'MRR'
        }
        
        for old, new in replacements.items():
            cleaned = cleaned.replace(old, new)
        
        return cleaned
    
    def create_title_slide(self, title, subtitle=None, date=None):
        """Create a professional title slide"""
        # Use branded generator if available
        if self.use_branding:
            return self.branded_generator.create_title_slide(title, subtitle, date)
        
        # Create title slide
        slide_layout = self.prs.slide_layouts[0] if len(self.prs.slide_layouts) > 0 else self.prs.slide_layouts[-1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add background gradient
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        shape.fill.gradient()
        shape.fill.gradient_angle = 90
        shape.fill.gradient_stops[0].color.rgb = RGBColor(240, 248, 255)  # Alice Blue
        shape.fill.gradient_stops[1].color.rgb = RGBColor(255, 255, 255)  # White
        shape.line.fill.background()
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        title_frame = title_shape.text_frame
        title_frame.text = title
        
        # Style the title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(48)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.color.rgb = RGBColor(37, 64, 97)  # Dark blue
        
        # Add subtitle if provided
        if subtitle:
            subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.text = subtitle
            
            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.font.size = Pt(28)
            subtitle_paragraph.alignment = PP_ALIGN.CENTER
            subtitle_paragraph.font.color.rgb = RGBColor(79, 129, 189)  # Medium blue
        
        # Add date
        if not date:
            from datetime import datetime
            date = datetime.now().strftime("%B %Y")
        
        date_shape = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
        date_frame = date_shape.text_frame
        date_frame.text = date
        
        date_paragraph = date_frame.paragraphs[0]
        date_paragraph.font.size = Pt(18)
        date_paragraph.alignment = PP_ALIGN.CENTER
        date_paragraph.font.color.rgb = RGBColor(128, 128, 128)  # Gray
        
        return slide
    
    def create_thank_you_slide(self):
        """Create a professional thank you/questions slide"""
        # Use branded generator if available
        if self.use_branding:
            return self.branded_generator.create_thank_you_slide()
        
        # Create slide
        slide_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[-1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add background gradient (similar to title slide)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        shape.fill.gradient()
        shape.fill.gradient_angle = 90
        shape.fill.gradient_stops[0].color.rgb = RGBColor(240, 248, 255)  # Alice Blue
        shape.fill.gradient_stops[1].color.rgb = RGBColor(255, 255, 255)  # White
        shape.line.fill.background()
        
        # Add "Thank You" text
        thank_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        thank_frame = thank_shape.text_frame
        thank_frame.text = "Thank You"
        
        thank_paragraph = thank_frame.paragraphs[0]
        thank_paragraph.font.size = Pt(54)
        thank_paragraph.font.bold = True
        thank_paragraph.alignment = PP_ALIGN.CENTER
        thank_paragraph.font.color.rgb = RGBColor(37, 64, 97)  # Dark blue
        
        # Add "Questions?" text
        questions_shape = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        questions_frame = questions_shape.text_frame
        questions_frame.text = "Questions?"
        
        questions_paragraph = questions_frame.paragraphs[0]
        questions_paragraph.font.size = Pt(36)
        questions_paragraph.alignment = PP_ALIGN.CENTER
        questions_paragraph.font.color.rgb = RGBColor(79, 129, 189)  # Medium blue
        
        return slide
    
    def add_slide_numbers(self, start_from=1, exclude_first=True, exclude_last=True):
        """Add slide numbers to all slides"""
        total_slides = len(self.prs.slides)
        
        for idx, slide in enumerate(self.prs.slides):
            # Skip first slide (title) if requested
            if exclude_first and idx == 0:
                continue
            
            # Skip last slide (thank you) if requested
            if exclude_last and idx == total_slides - 1:
                continue
            
            # Add slide number
            slide_num = idx if not exclude_first else idx
            footer_shape = slide.shapes.add_textbox(
                Inches(8.5), Inches(7), Inches(1), Inches(0.3)
            )
            footer_frame = footer_shape.text_frame
            footer_frame.text = f"{slide_num}/{total_slides - (1 if exclude_first else 0) - (1 if exclude_last else 0)}"
            
            footer_paragraph = footer_frame.paragraphs[0]
            footer_paragraph.font.size = Pt(10)
            footer_paragraph.alignment = PP_ALIGN.RIGHT
            footer_paragraph.font.color.rgb = RGBColor(128, 128, 128)  # Gray
    
    def save_presentation(self, output_path):
        """Save the presentation to specified path"""
        # Use branded generator if available
        if self.use_branding:
            return self.branded_generator.save_presentation(output_path)
        
        # Fallback to original implementation
        self.prs.save(output_path)
        return output_path