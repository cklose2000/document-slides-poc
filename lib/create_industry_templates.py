"""
Create industry-specific PowerPoint templates with custom styling
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import os
import json

def create_industry_template(template_id: str, metadata_path: str, output_path: str):
    """Create an industry-specific PowerPoint template"""
    
    # Load metadata
    with open(metadata_path, 'r') as f:
        metadata = json.load(f)
    
    # Create new presentation
    prs = Presentation()
    
    # Get color scheme from metadata
    colors = metadata.get('colors', {})
    fonts = metadata.get('fonts', {})
    
    # Create title slide layout (slide 0)
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # Customize title slide
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    # Set title
    title.text = f"{metadata['name']} Template"
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.name = fonts.get('title_font', 'Arial')
    title_paragraph.font.size = Pt(fonts.get('title_size', 44))
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor.from_string(colors.get('primary', '#4F81BD').replace('#', ''))
    
    # Set subtitle
    subtitle.text = metadata.get('description', 'Professional presentation template')
    subtitle_paragraph = subtitle.text_frame.paragraphs[0]
    subtitle_paragraph.font.name = fonts.get('body_font', 'Arial')
    subtitle_paragraph.font.size = Pt(fonts.get('subtitle_size', 32))
    subtitle_paragraph.font.color.rgb = RGBColor.from_string(colors.get('text', '#333333').replace('#', ''))
    
    # Create content slide with title and content (slide 1)
    content_slide_layout = prs.slide_layouts[1]
    content_slide = prs.slides.add_slide(content_slide_layout)
    
    # Customize content slide
    content_title = content_slide.shapes.title
    content_title.text = "Key Metrics Dashboard"
    content_title_paragraph = content_title.text_frame.paragraphs[0]
    content_title_paragraph.font.name = fonts.get('title_font', 'Arial')
    content_title_paragraph.font.size = Pt(fonts.get('title_size', 44))
    content_title_paragraph.font.bold = True
    content_title_paragraph.font.color.rgb = RGBColor.from_string(colors.get('primary', '#4F81BD').replace('#', ''))
    
    # Add industry-specific content placeholder
    content_placeholder = content_slide.placeholders[1]
    content_text = content_placeholder.text_frame
    content_text.text = f"• Industry Focus: {metadata['name']}\n"
    
    # Add industry-specific metrics
    if 'industry_focus' in metadata and 'metrics' in metadata['industry_focus']:
        content_text.text += "• Key Metrics:\n"
        for metric in metadata['industry_focus']['metrics'][:5]:  # First 5 metrics
            content_text.text += f"  - {metric}\n"
    
    # Style the content
    for paragraph in content_text.paragraphs:
        paragraph.font.name = fonts.get('body_font', 'Arial')
        paragraph.font.size = Pt(fonts.get('body_size', 16))
        paragraph.font.color.rgb = RGBColor.from_string(colors.get('text', '#333333').replace('#', ''))
    
    # Create section header slide (slide 2)
    section_slide_layout = prs.slide_layouts[2]
    section_slide = prs.slides.add_slide(section_slide_layout)
    
    section_title = section_slide.shapes.title
    section_title.text = "Financial Performance"
    section_title_paragraph = section_title.text_frame.paragraphs[0]
    section_title_paragraph.font.name = fonts.get('title_font', 'Arial')
    section_title_paragraph.font.size = Pt(fonts.get('title_size', 44))
    section_title_paragraph.font.bold = True
    section_title_paragraph.font.color.rgb = RGBColor.from_string(colors.get('primary', '#4F81BD').replace('#', ''))
    section_title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Create two content slide (slide 3)
    two_content_layout = prs.slide_layouts[3]
    two_content_slide = prs.slides.add_slide(two_content_layout)
    
    two_content_title = two_content_slide.shapes.title
    two_content_title.text = "Analysis & Charts"
    two_content_title_paragraph = two_content_title.text_frame.paragraphs[0]
    two_content_title_paragraph.font.name = fonts.get('title_font', 'Arial')
    two_content_title_paragraph.font.size = Pt(fonts.get('title_size', 44))
    two_content_title_paragraph.font.bold = True
    two_content_title_paragraph.font.color.rgb = RGBColor.from_string(colors.get('primary', '#4F81BD').replace('#', ''))
    
    # Customize placeholders if they exist
    if len(two_content_slide.placeholders) > 1:
        left_placeholder = two_content_slide.placeholders[1]
        left_text = left_placeholder.text_frame
        left_text.text = "Key Insights:\n• Data Analysis\n• Trend Identification\n• Performance Metrics"
        
        for paragraph in left_text.paragraphs:
            paragraph.font.name = fonts.get('body_font', 'Arial')
            paragraph.font.size = Pt(fonts.get('body_size', 16))
            paragraph.font.color.rgb = RGBColor.from_string(colors.get('text', '#333333').replace('#', ''))
    
    if len(two_content_slide.placeholders) > 2:
        right_placeholder = two_content_slide.placeholders[2]
        right_text = right_placeholder.text_frame
        right_text.text = "[Chart Placeholder]\n\nThis area will contain auto-generated charts based on your data."
        
        for paragraph in right_text.paragraphs:
            paragraph.font.name = fonts.get('body_font', 'Arial')
            paragraph.font.size = Pt(fonts.get('body_size', 16))
            paragraph.font.color.rgb = RGBColor.from_string(colors.get('text', '#333333').replace('#', ''))
    
    # Create blank slide for custom content (slide 4)
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    blank_slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add a custom title to the blank slide
    title_shape = blank_slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Custom Analysis"
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = fonts.get('title_font', 'Arial')
    title_paragraph.font.size = Pt(fonts.get('title_size', 44))
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor.from_string(colors.get('primary', '#4F81BD').replace('#', ''))
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add footer with industry branding
    footer_shape = blank_slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    footer_frame = footer_shape.text_frame
    footer_frame.text = f"{metadata['name']} | {metadata.get('description', '')}"
    footer_paragraph = footer_frame.paragraphs[0]
    footer_paragraph.font.name = fonts.get('body_font', 'Arial')
    footer_paragraph.font.size = Pt(10)
    footer_paragraph.font.color.rgb = RGBColor.from_string(colors.get('neutral', '#666666').replace('#', ''))
    footer_paragraph.alignment = PP_ALIGN.CENTER
    
    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    # Save template
    prs.save(output_path)
    print(f"Industry template created: {output_path} for {metadata['name']}")
    
    return output_path

def create_all_industry_templates():
    """Create all industry-specific templates"""
    base_path = "/mnt/c/Users/cklos/document-slides-poc/templates"
    
    industries = ['saas', 'manufacturing', 'healthcare']
    
    for industry in industries:
        metadata_path = os.path.join(base_path, industry, 'metadata.json')
        output_path = os.path.join(base_path, industry, 'template.pptx')
        
        if os.path.exists(metadata_path):
            try:
                create_industry_template(industry, metadata_path, output_path)
                print(f"✅ Created {industry} template successfully")
            except Exception as e:
                print(f"❌ Failed to create {industry} template: {str(e)}")
        else:
            print(f"❌ Metadata not found for {industry}: {metadata_path}")

if __name__ == "__main__":
    create_all_industry_templates()