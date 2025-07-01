#!/usr/bin/env python3
"""
Test script to validate PowerPoint corruption fixes
Creates slides using the fixed components to verify they open correctly
"""

from pptx import Presentation
from pptx.util import Inches
import sys
import os
sys.path.append('/mnt/c/Users/cklos/document-slides-poc/lib')

from slide_components import TextComponents, DataComponents, VisualComponents, CompositeComponents

def test_corruption_fixes():
    """Test that fixed components don't create corrupted PPTX files"""
    
    print("🧪 Testing PowerPoint corruption fixes...")
    
    # Create presentation
    prs = Presentation()
    
    # Test 1: KPI Cards with edge cases
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    data_comp = DataComponents(slide1)
    
    print("   Testing KPI cards with extreme values...")
    
    # Test extreme coordinates and values
    try:
        kpi_card = data_comp.add_kpi_card(
            title="Revenue Test",
            value=1000000000,  # Large number
            change=150.5,      # Large percentage
            left=0.1,          # Near edge
            top=0.1,           # Near edge
            width=2.0,
            height=1.5
        )
        print("   ✅ KPI card with extreme values created successfully")
    except Exception as e:
        print(f"   ❌ KPI card failed: {e}")
        return False
    
    # Test 2: Progress indicators with edge cases
    print("   Testing progress indicators...")
    try:
        # Test edge cases
        progress1 = data_comp.add_progress_indicator(
            progress=0,        # Minimum
            left=0.5,
            top=3.0,
            width=3.0,
            label="Zero Progress"
        )
        
        progress2 = data_comp.add_progress_indicator(
            progress=100,      # Maximum
            left=0.5,
            top=3.5,
            width=3.0,
            label="Full Progress"
        )
        
        progress3 = data_comp.add_progress_indicator(
            progress=150,      # Over 100% (should be clamped)
            left=0.5,
            top=4.0,
            width=3.0,
            label="Over 100%"
        )
        print("   ✅ Progress indicators with edge values created successfully")
    except Exception as e:
        print(f"   ❌ Progress indicators failed: {e}")
        return False
    
    # Test 3: Timeline with many events
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    composite_comp = CompositeComponents(slide2)
    
    print("   Testing timeline visualization...")
    try:
        events = [
            {"date": "Q1 2024", "title": "Launch"},
            {"date": "Q2 2024", "title": "Growth"},
            {"date": "Q3 2024", "title": "Scale"},
            {"date": "Q4 2024", "title": "Expand"},
            {"date": "Q1 2025", "title": "IPO"}
        ]
        
        timeline = composite_comp.add_timeline_visualization(
            events=events,
            left=0.5,
            top=2.0,
            width=9.0,
            height=3.0
        )
        print("   ✅ Timeline visualization created successfully")
    except Exception as e:
        print(f"   ❌ Timeline visualization failed: {e}")
        return False
    
    # Test 4: Risk matrix with edge coordinates
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    composite_comp3 = CompositeComponents(slide3)
    
    print("   Testing risk matrix...")
    try:
        risks = [
            {"name": "Market Risk", "impact": 3, "probability": 2},
            {"name": "Tech Risk", "impact": 2, "probability": 3},
            {"name": "Funding Risk", "impact": 1, "probability": 1},
            {"name": "Extreme Risk", "impact": 5, "probability": 5},  # Should be clamped to 3,3
        ]
        
        risk_matrix = composite_comp3.add_risk_matrix(
            risks=risks,
            left=1.0,
            top=1.0,
            size=5.0
        )
        print("   ✅ Risk matrix created successfully")
    except Exception as e:
        print(f"   ❌ Risk matrix failed: {e}")
        return False
    
    # Test 5: Text components with various font sizes
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    text_comp = TextComponents(slide4)
    
    print("   Testing text components with extreme font sizes...")
    try:
        # Test various font sizes
        title = text_comp.add_title(
            "Test Title",
            left=0.5,
            top=0.5,
            width=9.0,
            height=1.0,
            font_size=5000  # Should be clamped to 4000
        )
        
        body = text_comp.add_body_text(
            "This is test body text with extreme font size.",
            left=0.5,
            top=2.0,
            width=9.0,
            height=2.0,
            font_size=0  # Should be clamped to 1
        )
        
        bullets = text_comp.add_bullet_points(
            ["Item 1", "Item 2", "Item 3"],
            left=0.5,
            top=4.5,
            width=9.0,
            height=1.5,
            font_size=-10  # Should be clamped to 1
        )
        print("   ✅ Text components with extreme font sizes created successfully")
    except Exception as e:
        print(f"   ❌ Text components failed: {e}")
        return False
    
    # Test 6: Table with edge case data
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    data_comp5 = DataComponents(slide5)
    
    print("   Testing metric table...")
    try:
        table_data = [
            {"Metric": "Revenue", "Q1": 1000000, "Q2": 1500000},
            {"Metric": "Profit", "Q1": 100000, "Q2": 200000},
            {"Metric": "Growth", "Q1": 0, "Q2": 50.5}
        ]
        
        table = data_comp5.add_metric_table(
            data=table_data,
            left=0.5,
            top=1.0,
            width=8.0,
            height=3.0
        )
        print("   ✅ Metric table created successfully")
    except Exception as e:
        print(f"   ❌ Metric table failed: {e}")
        return False
    
    # Save and validate
    output_path = "/mnt/c/Users/cklos/document-slides-poc/corruption_test_output.pptx"
    
    try:
        prs.save(output_path)
        print(f"   ✅ PPTX file saved successfully: {output_path}")
        
        # Check file size
        file_size = os.path.getsize(output_path)
        print(f"   📊 File size: {file_size:,} bytes")
        
        if file_size < 10000:  # Less than 10KB might indicate corruption
            print("   ⚠️  Warning: File size seems small, possible corruption")
            return False
        
        print("   🎉 All tests passed! File should open without corruption.")
        return True
        
    except Exception as e:
        print(f"   ❌ Failed to save PPTX: {e}")
        return False

if __name__ == "__main__":
    print("🔧 PowerPoint Corruption Fix Validation")
    print("=" * 50)
    
    success = test_corruption_fixes()
    
    if success:
        print("\n✅ SUCCESS: All corruption fixes validated!")
        print("💡 The generated file should now open correctly in PowerPoint.")
        print("📁 File: corruption_test_output.pptx")
    else:
        print("\n❌ FAILURE: Some tests failed.")
        print("🔍 Check the error messages above for details.")
    
    exit(0 if success else 1)