"""
Quick script to set up demo data and test the POC
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_sample_template():
    """Create a basic firm template"""
    template_dir = 'templates'
    if not os.path.exists(template_dir):
        os.makedirs(template_dir)
    
    template_path = os.path.join(template_dir, 'firm_template.pptx')
    
    # Create a new presentation
    prs = Presentation()
    
    # Customize the master slide layouts
    # Layout 0: Title slide
    title_slide_layout = prs.slide_layouts[0]
    
    # Layout 1: Title and content
    content_slide_layout = prs.slide_layouts[1]
    
    # Add a sample title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Company Analysis"
    subtitle.text = "Generated from Data Room Documents"
    
    # Style the title slide
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)  # Dark blue
    
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 114, 196)  # Lighter blue
    
    # Remove the sample slide (we just wanted to set up the template)
    rIndex = len(prs.slides) - 1
    del prs.slides._sldIdLst[rIndex]
    
    # Save the template
    prs.save(template_path)
    print(f"✅ Template created at: {template_path}")
    return template_path

def create_sample_data():
    """Create sample data files for testing"""
    # This would create sample Excel, PDF, and Word files
    # For now, just create placeholder instructions
    
    sample_dir = 'sample_data'
    if not os.path.exists(sample_dir):
        os.makedirs(sample_dir)
    
    # Create a sample instruction file
    instructions = """
Sample Data for Testing
======================

To test the POC, create these sample files:

1. sample_financials.xlsx:
   - Sheet 'Summary' with revenue, profit, growth metrics
   - Use formulas like =SUM(B1:B10) for calculations
   - Put company name in cell A1

2. sample_memo.pdf:
   - Company overview document
   - Include industry, business model, market position
   - Add performance highlights section

3. sample_business_plan.docx:
   - Executive summary
   - Market analysis
   - Financial projections
   - Use proper heading styles

Upload these files through the web interface to test extraction and slide generation.
"""
    
    with open(os.path.join(sample_dir, 'README.txt'), 'w') as f:
        f.write(instructions)
    
    print(f"✅ Sample data instructions created in: {sample_dir}/README.txt")

def create_run_script():
    """Create a simple run script"""
    run_script = """#!/bin/bash
# Document to Slides POC - Run Script

echo "🚀 Starting Document to Slides POC"
echo "=================================="

# Check if virtual environment exists
if [ ! -d "venv" ]; then
    echo "📦 Creating virtual environment..."
    python -m venv venv
fi

# Activate virtual environment
source venv/bin/activate

# Install requirements
echo "📚 Installing dependencies..."
pip install -r requirements.txt

# Set environment variables (update these with your actual values)
export DATABASE_URL="sqlite:///app.db"  # or your database URL
export JWT_SECRET="your-jwt-secret-here"
# export OPENAI_API_KEY="your-openai-key-here"  # Uncomment and add your key
# export LLMWHISPERER_API_KEY="your-llm-whisperer-key"  # For PDF processing

echo "🌐 Starting server..."
echo "📋 Web interface will be available at: http://localhost:5000/static/presentation.html"
echo "🔍 Health check: http://localhost:5000/health"
echo "📊 Preview API: http://localhost:5000/api/generate-slides/preview"

cd api
python generate_slides.py
"""
    
    with open('run.sh', 'w') as f:
        f.write(run_script)
    
    # Make it executable
    os.chmod('run.sh', 0o755)
    
    print("✅ Run script created: ./run.sh")

def create_readme():
    """Create a comprehensive README"""
    readme_content = """# Document to Slides POC

Transform data room documents (PDF, Excel, Word) into branded PowerPoint presentations with source attribution.

## Features

- 📄 **Multi-format Support**: PDF, Excel (.xlsx), Word (.docx)
- 🧠 **AI Analysis**: Extract key metrics and insights using LLM
- 📊 **Source Attribution**: Track data back to specific cells/pages
- 🎯 **Professional Slides**: Generate branded PowerPoint presentations
- 🔍 **Preview Mode**: Test extraction without generating slides

## Quick Start

1. **Install Dependencies**:
   ```bash
   pip install -r requirements.txt
   ```

2. **Set Environment Variables**:
   ```bash
   export OPENAI_API_KEY="your-key-here"  # Optional, for AI analysis
   ```

3. **Run the Server**:
   ```bash
   chmod +x run.sh
   ./run.sh
   ```

4. **Open Web Interface**:
   Visit: http://localhost:5000/static/presentation.html

## File Structure

```
document-slides-poc/
├── lib/
│   ├── excel_extractor.py      # Excel parsing with cell references
│   ├── word_extractor.py       # Word document structure extraction
│   ├── slide_generator.py      # PowerPoint generation
│   └── llm_slides.py          # AI document analysis
├── api/
│   └── generate_slides.py     # Flask API endpoints
├── static/
│   └── presentation.html      # Web interface
├── templates/
│   └── firm_template.pptx     # PowerPoint template
├── requirements.txt           # Python dependencies
└── demo_setup.py             # Setup script
```

## API Endpoints

- `POST /api/generate-slides` - Generate PowerPoint from documents
- `POST /api/generate-slides/preview` - Preview extraction results
- `GET /health` - Health check

## Usage

1. **Upload Documents**: Drag & drop or select PDF/Excel/Word files
2. **Preview Extraction**: Click "Preview Extraction" to see parsed data
3. **Generate Slides**: Click "Generate Slides" to create PowerPoint
4. **Download**: Save the generated presentation

## Sample Slide Types

- **Financial Summary**: Key metrics with source attribution
- **Company Overview**: Business description and industry
- **Key Insights**: AI-extracted highlights

## Environment Variables

- `OPENAI_API_KEY`: For AI document analysis (optional)
- `DATABASE_URL`: Database connection (for future features)
- `JWT_SECRET`: Security key (for future auth features)

## Source Attribution

Each data point includes:
- Document name
- Sheet/page reference
- Specific cell/section location
- Confidence score (when using AI)

## Testing

Use the sample data instructions in `sample_data/README.txt` to create test files.

## Next Steps

- Add more slide templates
- Implement user authentication
- Add database storage for analysis results
- Enhance AI prompts for better extraction
- Add support for more file formats

---

**Note**: This is a POC demonstrating core functionality. Production deployment would require additional security, error handling, and scalability considerations.
"""
    
    with open('README.md', 'w') as f:
        f.write(readme_content)
    
    print("✅ README.md created")

def main():
    """Set up the demo environment"""
    print("🔧 Setting up Document to Slides POC Demo")
    print("=" * 50)
    
    # Create template
    create_sample_template()
    
    # Create sample data instructions
    create_sample_data()
    
    # Create run script
    create_run_script()
    
    # Create README
    create_readme()
    
    print("\n🎉 Demo setup complete!")
    print("\nNext steps:")
    print("1. Install dependencies: pip install -r requirements.txt")
    print("2. Set your OpenAI API key (optional): export OPENAI_API_KEY='your-key'")
    print("3. Run the server: ./run.sh")
    print("4. Open http://localhost:5000/static/presentation.html")
    print("5. Upload test documents and generate slides!")

if __name__ == "__main__":
    main()